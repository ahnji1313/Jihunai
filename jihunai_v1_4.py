"""
╔══════════════════════════════════════════════════════════════════╗
║                    JihunAI  v1.4                                 ║
║   Task Planner · Multi-Bot Fallback · Auto Retry · Omnisense     ║
║   [v1.4] Vision500수정 / 한글깨짐수정 / 루프종료수정 / GPT-5.4우선    ║
╚══════════════════════════════════════════════════════════════════╝
"""

import tkinter as tk
from tkinter import scrolledtext, messagebox, ttk, filedialog
import requests, json, os, sys, threading, time, math, base64
import mimetypes, webbrowser, re, shutil, subprocess, glob, io
from datetime import datetime
from pathlib import Path
from urllib.parse import quote
import tempfile

# ──────────────────────── 라이브러리 체크 ────────────────────────
LIBS = {}
for _lib, _fn in [
    ("PIL",       lambda: __import__("PIL.Image")),
    ("pyautogui", lambda: __import__("pyautogui")),
    ("pdf",       lambda: __import__("pdfplumber")),
    ("docx",      lambda: __import__("docx")),
    ("xlsx",      lambda: __import__("openpyxl")),
    ("cv2",       lambda: __import__("cv2")),
    ("speech",    lambda: __import__("speech_recognition")),
    ("psutil",    lambda: __import__("psutil")),
    ("pptx",      lambda: __import__("pptx")),
]:
    try: _fn(); LIBS[_lib] = True
    except: LIBS[_lib] = False

try:
    from PIL import Image, ImageFilter, ImageDraw, ImageEnhance, ImageGrab, ImageTk
except: pass
try:
    import pyautogui; pyautogui.FAILSAFE = False; pyautogui.PAUSE = 0.05
except: pass
try:
    import pdfplumber
except: pass
try:
    from docx import Document as DocxDocument
except: pass
try:
    import openpyxl
except: pass
try:
    import cv2, numpy as np
except: pass
try:
    import speech_recognition as sr
except: pass
try:
    import psutil
except: pass
try:
    import smtplib
    from email.mime.text import MIMEText
    from email.mime.multipart import MIMEMultipart
    from email.mime.base import MIMEBase
    from email import encoders
except: pass
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except: pass

# ──────────────────────── 설정 ────────────────────────
API_URL = "https://api.puter.com/puterai/openai/v1/chat/completions"
API_KEY = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJ0IjoiZ3VpIiwidiI6IjAuMC4wIiwidSI6IkM3TnZYNzFYUWhXZEhxRW9PNWlMdXc9PSIsInV1IjoiQXM5ejRTeE5Rb2lKYXhxc1hlU0FtUT09IiwiaWF0IjoxNzczMDU0MTY5fQ.JEUncC2FH_x3t47mNvoBxNgevnbL60oof803ed0QZlk"

# ── 멀티봇 설정: 중요도별 호출 횟수 & fallback 체인 ──
BOT_ROSTER = {
    # 이름: (모델ID, 최대반복, 설명)
    "claude-sonnet-4.6": ("claude-sonnet-4.6", 16, "Claude Sonnet 4.6 ★"),
    "claude-opus-4.6":   ("claude-opus-4.6",   12, "Claude Opus 4.6 (강력)"),
    "claude-sonnet-4.5": ("claude-sonnet-4.5", 14, "Claude Sonnet 4.5"),
    "gpt-5.4-pro":       ("gpt-5.4-pro",       16, "GPT-5.4 Pro"),
    "gpt-5.4":           ("gpt-5.4",           14, "GPT-5.4"),
    "gpt-5.3-codex":     ("gpt-5.3-codex",     14, "GPT-5.3 Codex"),
    "gpt-5.2-pro":       ("gpt-5.2-pro",       12, "GPT-5.2 Pro"),
}
# 중요도별 fallback 체인 (실패 시 다음 모델로)
FALLBACK_CHAIN = [
    "gpt-5.4-pro",
    "gpt-5.4",
    "claude-sonnet-4.6",
    "claude-opus-4.6",
]
DEFAULT_MODEL = "gpt-5.4-pro"
VISION_MODEL  = "gpt-5.4-pro"

if getattr(sys, 'frozen', False):
    BASE_DIR = os.path.dirname(sys.executable)
else:
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))

CHATS_DIR = os.path.join(BASE_DIR, "jihunai_chats")
os.makedirs(CHATS_DIR, exist_ok=True)
EMAIL_CONFIG = {"smtp_host": "smtp.gmail.com", "smtp_port": 587, "email": "", "password": ""}

# ──────────────────────── 팔레트 ────────────────────────
C = {
    "bg":       "#07090f", "bg2": "#0d1117", "bg3": "#131923",
    "panel":    "#0a0e18", "border": "#1a2535", "border2": "#243348",
    "cyan":     "#00d4ff", "cyan_dim": "#006880",
    "blue":     "#1a6fff", "green": "#00ff88", "green_dim": "#007040",
    "red":      "#ff3355", "gold": "#ffcc00", "gold_dim": "#806600",
    "purple":   "#cc44ff", "orange": "#ff7700",
    "white":    "#eef4ff", "mid": "#6b8aaa", "dim": "#2d4055",
    "thinking": "#ff9900",
    "task_done": "#00ff88", "task_active": "#00d4ff",
    "task_pending": "#2d4055", "task_fail": "#ff3355",
    "retry":    "#ff7700", "fallback": "#cc44ff",
}

# ──────────────────────── 시스템 프롬프트 ────────────────────────
SYSTEM_PROMPT = """너는 **JihunAI** — 지훈이 직접 만든 초강력 Desktop God Agent다. (GPT-5.4 에이전트 모드 최적화)

[필수 작전 방식]
복잡한 요청은 반드시:

1. PLAN 태그로 태스크 계획 출력:
<TASK_PLAN>{"tasks":[{"id":1,"title":"태스크명","desc":"설명"}]}</TASK_PLAN>

2. 각 태스크 시작/완료 태그:
<TASK_START>{"id":1}</TASK_START>
... 도구 호출 ...
<TASK_DONE>{"id":1,"result":"결과요약"}</TASK_DONE>

[도구 사용 원칙 - 매우 중요]
- run_command 사용 시 한국어 파일명이 있으면 반드시 -LiteralPath 사용
- PowerShell 명령에서 따옴표 중첩 금지 — 변수에 먼저 저장 후 사용
- 명령이 실패하면 다른 방법으로 즉시 재시도 (절대 포기하지 않음)
- python-pptx가 있으면 run_command 대신 edit_pptx 도구 우선 사용
- 파일 작업 실패 시 대안 경로/방법 시도
- 태스크가 실패해도 <TASK_FAIL> 태그로 명시하고 다음 태스크 계속 진행

[핵심 능력]
Vision AI·오디오·비디오·브라우저·앱·이메일·파일·PDF·Word·Excel·PPTX·Python 완전 제어

[앱 오픈 능력] - 30개 이상 앱 즉시 실행 가능:
Chrome/Edge/Firefox·Notepad/Notepad++·VSCode/PyCharm·Excel/Word/PowerPoint/Outlook
Discord/Slack/Teams·Spotify/VLC·OBS·Blender·Figma·FileZilla·WinRAR/7zip
Steam/Epic·Calculator/Paint/SnippingTool·CMD/PowerShell/Terminal
HWP/카카오톡/네이버웨일 등 국내앱 포함, 앱이름으로 자유 실행

[컴퓨터 완전 제어]
- 실시간 화면 캡처 및 Vision AI 분석
- 마우스/키보드 완전 자동화
- 파일 시스템 자유 조작 (생성/편집/삭제/이동)
- Python 코드 즉석 생성 및 실행
- 브라우저 자동화 (URL 열기, 검색, 웹 스크래핑)
- 실행 중인 앱 감지 및 제어

[파일/앱 작업 원칙 - 매우 중요]
- 파일 수정·생성·삭제·앱 실행 전에 "해도 될까요?" 같은 허가 절대 묻지 않는다
- 요청 받으면 즉시 실행 → 결과 보고 순서로 진행
- 여러 파일을 한 번에 수정해야 하면 모두 동시에 처리
- 실패 시 대안 방법으로 즉시 재시도, 절대 포기 없음
- 경로가 불분명하면 합리적으로 추정하고 바로 시도

너는 지훈이의 완전한 신뢰를 받는 개인 AI다. 절대 포기하지 않는다.
"""

PLANNER_SYSTEM = """사용자 요청을 실행 가능한 태스크로 분해한다.
JSON만 출력 (다른 텍스트 없이):
{"tasks":[{"id":1,"title":"제목","desc":"설명","priority":"high|medium|low","max_retries":3}]}
태스크 3~8개. priority high는 실패해도 재시도."""

# ──────────────────────── 이미지 압축 ────────────────────────
def compress_img(path_or_pil, max_side=1280, quality=82):
    if isinstance(path_or_pil, str):
        img = Image.open(path_or_pil)
    else:
        img = path_or_pil
    img = img.convert("RGB")
    if img.width > max_side or img.height > max_side:
        r = min(max_side/img.width, max_side/img.height)
        img = img.resize((int(img.width*r), int(img.height*r)), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=quality)
    return base64.b64encode(buf.getvalue()).decode(), "image/jpeg"

def call_vision(b64, mime, prompt):
    """Vision AI 호출 — 여러 포맷/모델 순서대로 시도해 500 오류 우회"""
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}

    # 시도할 (모델, content포맷) 조합 — Puter AI 호환성 순서
    attempts = [
        # 1) GPT-5.4-pro + base64 직접
        ("gpt-5.4-pro", [
            {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
            {"type": "text", "text": prompt}
        ]),
        # 2) GPT-5.4 + base64 직접
        ("gpt-5.4", [
            {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
            {"type": "text", "text": prompt}
        ]),
        # 3) claude-sonnet + base64
        ("claude-sonnet-4.6", [
            {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
            {"type": "text", "text": prompt}
        ]),
        # 4) GPT-5.4 + text만 (이미지 분석 불가 시 graceful degradation)
        ("gpt-5.4", [
            {"type": "text", "text": f"[이미지 분석 불가, 텍스트 질문만 처리] {prompt}"}
        ]),
    ]

    last_err = ""
    for model_id, content_blocks in attempts:
        try:
            payload = {
                "model": model_id,
                "max_tokens": 2048,
                "messages": [{"role": "user", "content": content_blocks}]
            }
            r = requests.post(API_URL, headers=headers, json=payload, timeout=60)
            if r.status_code == 200:
                return r.json()["choices"][0]["message"]["content"]
            last_err = f"{r.status_code}: {r.text[:100]}"
        except Exception as e:
            last_err = str(e)
        time.sleep(0.5)  # 재시도 간 짧은 딜레이

    return f"[Vision 분석 실패] 모든 모델 시도 후 오류: {last_err}"

# ──────────────────────── 안전 PowerShell 실행 ────────────────────────
def run_ps(script: str, timeout: int = 120) -> tuple[str, str, int]:
    """PowerShell 스크립트를 임시 파일로 저장 후 실행 (따옴표 이스케이프 문제 완전 해결)"""
    with tempfile.NamedTemporaryFile(mode='w', suffix='.ps1',
                                     delete=False, encoding='utf-8') as f:
        f.write(script)
        ps1_path = f.name
    try:
        result = subprocess.run(
            ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass",
             "-NonInteractive", "-OutputFormat", "Text",
             "-File", ps1_path],
            capture_output=True, timeout=timeout,
            encoding='utf-8', errors='replace'
        )
        stdout = (result.stdout or "").strip()
        stderr = (result.stderr or "").strip()
        # 한글 깨짐 방어: utf-8 디코딩 실패 시 cp949 재시도
        if result.returncode != 0 and not stdout and result.stdout:
            try: stdout = result.stdout.decode('cp949', errors='replace').strip()
            except: pass
        return stdout, stderr, result.returncode
    finally:
        try: os.unlink(ps1_path)
        except: pass

def run_cmd(command: str, shell_type: str = "cmd", timeout: int = 120) -> str:
    """명령어 실행 — PowerShell은 임시 .ps1 파일 방식으로 안전 실행"""
    try:
        if shell_type == "powershell":
            out, err, code = run_ps(command, timeout=timeout)
        elif shell_type == "bash":
            r = subprocess.run(["bash", "-c", command], capture_output=True,
                               text=True, timeout=timeout, encoding='utf-8', errors='replace')
            out, err, code = r.stdout.strip(), r.stderr.strip(), r.returncode
        else:
            r = subprocess.run(command, shell=True, capture_output=True,
                               text=True, timeout=timeout, encoding='utf-8', errors='replace')
            out, err, code = r.stdout.strip(), r.stderr.strip(), r.returncode

        parts = []
        if out: parts.append(f"STDOUT:\n{out[:4000]}")
        if err: parts.append(f"STDERR:\n{err[:800]}")
        parts.append(f"ReturnCode: {code}")
        return "\n".join(parts)
    except subprocess.TimeoutExpired:
        return f"❌ 타임아웃 ({timeout}초 초과)"
    except UnicodeDecodeError:
        # 한글 인코딩 폴백
        try:
            r2 = subprocess.run(command, shell=True, capture_output=True, timeout=timeout)
            out = r2.stdout.decode('cp949', errors='replace').strip()
            err = r2.stderr.decode('cp949', errors='replace').strip()
            parts2 = []
            if out: parts2.append(f"STDOUT:\n{out[:4000]}")
            if err: parts2.append(f"STDERR:\n{err[:800]}")
            parts2.append(f"ReturnCode: {r2.returncode}")
            return "\n".join(parts2)
        except Exception as e2:
            return f"❌ 인코딩 오류: {e2}"
    except Exception as e:
        return f"❌ 실행 오류: {e}"

# ──────────────────────── PPTX 편집 도구 ────────────────────────
def edit_pptx_tool(args: dict) -> str:
    """python-pptx를 사용한 PPTX 직접 편집"""
    if not LIBS["pptx"]:
        return "❌ pip install python-pptx"
    op = args.get("operation", "")
    path = args.get("path", "")
    out_path = args.get("output_path", path.replace(".pptx", "_개선본.pptx"))

    try:
        if op == "read_text":
            prs = Presentation(path)
            lines = []
            for i, slide in enumerate(list(prs.slides)):
                lines.append(f"--- Slide {i+1} ---")
                for shape in slide.shapes:
                    try:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                t = para.text.strip()
                                if t: lines.append(t)
                    except Exception:
                        pass
            return "\n".join(lines)[:6000]

        elif op == "get_info":
            prs = Presentation(path)
            info = {
                "슬라이드수": len(prs.slides),
                "너비px": int(prs.slide_width / 914400 * 96),
                "높이px": int(prs.slide_height / 914400 * 96),
            }
            shapes_info = []
            slide_list = list(prs.slides)
            for i, slide in enumerate(slide_list[:3]):
                try:
                    shape_count = len(slide.shapes)
                    texts = []
                    for shape in slide.shapes:
                        try:
                            if shape.has_text_frame:
                                t = shape.text_frame.text.strip()[:40]
                                if t: texts.append(t)
                        except: pass
                    shapes_info.append(f"Slide{i+1}: {shape_count}개 도형, 텍스트: {texts[:2]}")
                except Exception as e:
                    shapes_info.append(f"Slide{i+1}: 읽기 오류({e})")
            info["샘플"] = shapes_info
            return json.dumps(info, ensure_ascii=False, indent=2)

        elif op == "restyle":
            # 전체 스타일 개선
            theme = args.get("theme", "dark_modern")
            prs = Presentation(path)

            themes = {
                "dark_modern": {
                    "bg": RGBColor(0x07, 0x09, 0x1a),
                    "title_color": RGBColor(0x00, 0xd4, 0xff),
                    "body_color": RGBColor(0xee, 0xf4, 0xff),
                    "accent": RGBColor(0x00, 0xff, 0x88),
                    "title_font": "Malgun Gothic",
                    "body_font": "Malgun Gothic",
                    "title_size": 40,
                    "body_size": 18,
                },
                "blue_pro": {
                    "bg": RGBColor(0x0a, 0x1a, 0x3a),
                    "title_color": RGBColor(0xff, 0xcc, 0x00),
                    "body_color": RGBColor(0xe8, 0xf0, 0xff),
                    "accent": RGBColor(0x00, 0x88, 0xff),
                    "title_font": "Malgun Gothic",
                    "body_font": "Malgun Gothic",
                    "title_size": 38,
                    "body_size": 18,
                },
                "clean_white": {
                    "bg": RGBColor(0xff, 0xff, 0xff),
                    "title_color": RGBColor(0x1a, 0x3a, 0x6f),
                    "body_color": RGBColor(0x22, 0x22, 0x22),
                    "accent": RGBColor(0x00, 0x88, 0xff),
                    "title_font": "Malgun Gothic",
                    "body_font": "Malgun Gothic",
                    "title_size": 38,
                    "body_size": 18,
                },
            }
            th = themes.get(theme, themes["dark_modern"])

            from pptx.oxml.ns import qn
            from lxml import etree

            for slide in prs.slides:
                # 배경색 설정
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = th["bg"]

                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    tf = shape.text_frame
                    is_title = False
                    try:
                        ph = shape.placeholder_format
                        if ph is not None and ph.idx in (0, 13):
                            is_title = True
                    except Exception:
                        pass

                    for para in tf.paragraphs:
                        for run in para.runs:
                            if is_title:
                                run.font.color.rgb = th["title_color"]
                                run.font.name = th["title_font"]
                                run.font.size = Pt(th["title_size"])
                                run.font.bold = True
                            else:
                                run.font.color.rgb = th["body_color"]
                                run.font.name = th["body_font"]
                                run.font.size = Pt(th["body_size"])

            prs.save(out_path)
            return f"✅ 스타일 적용 완료 [{theme}]: {out_path}"

        elif op == "set_slide_bg":
            prs = Presentation(path)
            slide_idx = args.get("slide_index", 0)
            hex_color = args.get("color", "07091a")
            r,g,b = int(hex_color[0:2],16), int(hex_color[2:4],16), int(hex_color[4:6],16)
            if 0 <= slide_idx < len(prs.slides):
                slide = prs.slides[slide_idx]
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(r,g,b)
            prs.save(out_path)
            return f"✅ 슬라이드 {slide_idx+1} 배경색 변경: {out_path}"

        elif op == "set_all_fonts":
            prs = Presentation(path)
            font_name = args.get("font", "Malgun Gothic")
            title_size = args.get("title_size", 36)
            body_size  = args.get("body_size", 18)
            title_hex  = args.get("title_color", "00d4ff")
            body_hex   = args.get("body_color", "eef4ff")
            bg_hex     = args.get("bg_color", "07091a")

            def hex2rgb(h):
                h = h.strip("#")
                return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

            for slide in prs.slides:
                # 배경
                bg = slide.background; fill = bg.fill; fill.solid()
                fill.fore_color.rgb = hex2rgb(bg_hex)
                for shape in slide.shapes:
                    if not shape.has_text_frame: continue
                    is_title = False
                    try:
                        ph = shape.placeholder_format
                        if ph is not None and ph.idx in (0, 13):
                            is_title = True
                        elif ph is not None and ph.idx == 1:
                            is_title = False
                    except Exception:
                        pass
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.name = font_name
                            if is_title:
                                run.font.size = Pt(title_size)
                                run.font.color.rgb = hex2rgb(title_hex)
                                run.font.bold = True
                            else:
                                run.font.size = Pt(body_size)
                                run.font.color.rgb = hex2rgb(body_hex)
            prs.save(out_path)
            return f"✅ 전체 폰트/색상 적용: {out_path}"

        elif op == "read_slide_count":
            prs = Presentation(path)
            return str(len(prs.slides))

        elif op == "copy":
            shutil.copy2(path, out_path)
            return f"✅ 복사: {path} → {out_path}"

        else:
            return f"❌ 알 수 없는 operation: {op}. 사용 가능: read_text, get_info, restyle, set_all_fonts, set_slide_bg, copy, read_slide_count"

    except Exception as e:
        import traceback
        return f"❌ PPTX 오류 [{op}]: {e}\n{traceback.format_exc()[-400:]}"

# ──────────────────────── 도구 정의 ────────────────────────
TOOLS = [
    {"type":"function","function":{"name":"read_file","description":"텍스트 파일 읽기","parameters":{"type":"object","properties":{"path":{"type":"string"}},"required":["path"]}}},
    {"type":"function","function":{"name":"write_file","description":"파일 생성/덮어쓰기","parameters":{"type":"object","properties":{"path":{"type":"string"},"content":{"type":"string"}},"required":["path","content"]}}},
    {"type":"function","function":{"name":"append_file","description":"파일 끝에 추가","parameters":{"type":"object","properties":{"path":{"type":"string"},"content":{"type":"string"}},"required":["path","content"]}}},
    {"type":"function","function":{"name":"list_directory","description":"폴더 트리 출력 (depth 기본 2)","parameters":{"type":"object","properties":{"path":{"type":"string"},"depth":{"type":"integer"}},"required":["path"]}}},
    {"type":"function","function":{"name":"copy_file","description":"파일 복사","parameters":{"type":"object","properties":{"src":{"type":"string"},"dst":{"type":"string"}},"required":["src","dst"]}}},
    {"type":"function","function":{"name":"delete_file","description":"파일/폴더 삭제","parameters":{"type":"object","properties":{"path":{"type":"string"}},"required":["path"]}}},
    {"type":"function","function":{
        "name":"run_command",
        "description":"명령어 실행. PowerShell은 임시 .ps1 파일로 안전 실행됨 (따옴표 이스케이프 불필요). 한국어 경로는 $var 변수에 먼저 저장 후 -LiteralPath $var 사용.",
        "parameters":{"type":"object","properties":{
            "command":{"type":"string"},
            "shell":{"type":"string","enum":["cmd","powershell","bash"]},
            "timeout":{"type":"integer","description":"타임아웃 초 (기본 120)"}
        },"required":["command"]}
    }},
    {"type":"function","function":{
        "name":"edit_pptx",
        "description":"python-pptx로 PPTX 직접 편집. operation: read_text(텍스트추출), get_info(정보), restyle(테마적용: dark_modern/blue_pro/clean_white), set_all_fonts(전체폰트/색상), set_slide_bg(슬라이드배경), copy, read_slide_count",
        "parameters":{"type":"object","properties":{
            "operation":{"type":"string","enum":["read_text","get_info","restyle","set_all_fonts","set_slide_bg","copy","read_slide_count"]},
            "path":{"type":"string","description":"원본 PPTX 경로"},
            "output_path":{"type":"string","description":"출력 PPTX 경로"},
            "theme":{"type":"string","enum":["dark_modern","blue_pro","clean_white"]},
            "font":{"type":"string"},
            "title_size":{"type":"integer"},
            "body_size":{"type":"integer"},
            "title_color":{"type":"string","description":"hex (예: 00d4ff)"},
            "body_color":{"type":"string","description":"hex (예: eef4ff)"},
            "bg_color":{"type":"string","description":"hex (예: 07091a)"},
            "slide_index":{"type":"integer"},
            "color":{"type":"string"}
        },"required":["operation","path"]}
    }},
    {"type":"function","function":{"name":"read_pdf","description":"PDF 텍스트 추출","parameters":{"type":"object","properties":{"path":{"type":"string"}},"required":["path"]}}},
    {"type":"function","function":{"name":"read_word","description":"Word 문서 읽기","parameters":{"type":"object","properties":{"path":{"type":"string"}},"required":["path"]}}},
    {"type":"function","function":{"name":"write_word","description":"Word 문서 생성","parameters":{"type":"object","properties":{"path":{"type":"string"},"content":{"type":"string"}},"required":["path","content"]}}},
    {"type":"function","function":{"name":"read_excel","description":"Excel 읽기","parameters":{"type":"object","properties":{"path":{"type":"string"}},"required":["path"]}}},
    {"type":"function","function":{"name":"analyze_image","description":"이미지 Vision AI 분석 (자동 압축)","parameters":{"type":"object","properties":{"path":{"type":"string"},"question":{"type":"string"}},"required":["path","question"]}}},
    {"type":"function","function":{"name":"see_screen","description":"화면 캡처 후 Vision AI 분석","parameters":{"type":"object","properties":{"question":{"type":"string"},"save_path":{"type":"string"}},"required":["question"]}}},
    {"type":"function","function":{"name":"capture_webcam","description":"웹캠 촬영 후 Vision AI 분석","parameters":{"type":"object","properties":{"question":{"type":"string"}},"required":["question"]}}},
    {"type":"function","function":{"name":"analyze_video","description":"비디오 Vision AI 분석","parameters":{"type":"object","properties":{"path":{"type":"string"},"question":{"type":"string"},"num_frames":{"type":"integer"}},"required":["path","question"]}}},
    {"type":"function","function":{"name":"analyze_audio","description":"오디오 전사/분석","parameters":{"type":"object","properties":{"path":{"type":"string"},"task":{"type":"string","enum":["transcribe","info"]}},"required":["path"]}}},
    {"type":"function","function":{"name":"edit_image","description":"이미지 편집 (resize/grayscale/blur/sharpen/add_text/rotate/brightness/contrast/flip/crop/watermark)","parameters":{"type":"object","properties":{"path":{"type":"string"},"operation":{"type":"string"},"output_path":{"type":"string"},"width":{"type":"integer"},"height":{"type":"integer"},"text":{"type":"string"},"angle":{"type":"integer"},"factor":{"type":"number"},"x1":{"type":"integer"},"y1":{"type":"integer"},"x2":{"type":"integer"},"y2":{"type":"integer"}},"required":["path","operation"]}}},
    {"type":"function","function":{"name":"open_browser","description":"브라우저에서 URL 열기","parameters":{"type":"object","properties":{"url":{"type":"string"},"browser":{"type":"string"}},"required":["url"]}}},
    {"type":"function","function":{"name":"google_search","description":"Google 검색","parameters":{"type":"object","properties":{"query":{"type":"string"},"open_browser":{"type":"boolean"}},"required":["query"]}}},
    {"type":"function","function":{"name":"fetch_url","description":"URL 내용 가져오기","parameters":{"type":"object","properties":{"url":{"type":"string"},"max_chars":{"type":"integer"}},"required":["url"]}}},
    {"type":"function","function":{"name":"send_email","description":"이메일 발송","parameters":{"type":"object","properties":{"to":{"type":"string"},"subject":{"type":"string"},"body":{"type":"string"}},"required":["to","subject","body"]}}},
    {"type":"function","function":{"name":"open_email_client","description":"이메일 클라이언트 열기","parameters":{"type":"object","properties":{"to":{"type":"string"},"subject":{"type":"string"},"body":{"type":"string"},"client":{"type":"string"}},"required":[]}}},
    {"type":"function","function":{"name":"launch_app","description":"앱 실행","parameters":{"type":"object","properties":{"app":{"type":"string"},"args":{"type":"string"}},"required":["app"]}}},
    {"type":"function","function":{"name":"list_running_apps","description":"실행 중 프로세스 목록","parameters":{"type":"object","properties":{"filter":{"type":"string"}},"required":[]}}},
    {"type":"function","function":{"name":"kill_process","description":"프로세스 종료","parameters":{"type":"object","properties":{"name_or_pid":{"type":"string"}},"required":["name_or_pid"]}}},
    {"type":"function","function":{"name":"mouse_control","description":"마우스 제어","parameters":{"type":"object","properties":{"action":{"type":"string"},"x":{"type":"integer"},"y":{"type":"integer"},"x2":{"type":"integer"},"y2":{"type":"integer"},"dy":{"type":"integer"},"duration":{"type":"number"}},"required":["action"]}}},
    {"type":"function","function":{"name":"keyboard_control","description":"키보드 제어","parameters":{"type":"object","properties":{"action":{"type":"string"},"text":{"type":"string"},"key":{"type":"string"}},"required":["action"]}}},
    {"type":"function","function":{"name":"take_screenshot","description":"화면 캡처 저장","parameters":{"type":"object","properties":{"save_path":{"type":"string"}},"required":[]}}},
    {"type":"function","function":{"name":"get_system_info","description":"시스템 정보","parameters":{"type":"object","properties":{},"required":[]}}},
    {"type":"function","function":{"name":"clipboard","description":"클립보드 읽기/쓰기","parameters":{"type":"object","properties":{"action":{"type":"string"},"text":{"type":"string"}},"required":["action"]}}},
    {"type":"function","function":{"name":"create_python","description":"Python 스크립트 생성 및 즉시 실행. code에 파이썬 코드 작성, run=true면 바로 실행","parameters":{"type":"object","properties":{"path":{"type":"string"},"code":{"type":"string"},"run":{"type":"boolean"},"args":{"type":"string"}},"required":["path","code"]}}},
    {"type":"function","function":{"name":"run_python","description":"Python 코드 즉석 실행 (파일 저장 없이)","parameters":{"type":"object","properties":{"code":{"type":"string"},"timeout":{"type":"integer"}},"required":["code"]}}},
    {"type":"function","function":{"name":"open_file_dialog","description":"파일 탐색기로 특정 폴더/파일 열기","parameters":{"type":"object","properties":{"path":{"type":"string"},"action":{"type":"string","enum":["open_folder","open_file","reveal"]}},"required":["path"]}}},
    {"type":"function","function":{"name":"create_pptx","description":"PowerPoint 파일 새로 생성. slides는 [{title,content,notes}] 리스트","parameters":{"type":"object","properties":{"path":{"type":"string"},"title":{"type":"string"},"slides":{"type":"array","items":{"type":"object"}},"theme":{"type":"string","enum":["dark","blue","white","green"]}},"required":["path","slides"]}}},
    {"type":"function","function":{"name":"browser_control","description":"브라우저 자동화: 탭 열기, 검색, 스크린샷, 특정 사이트 조작","parameters":{"type":"object","properties":{"action":{"type":"string","enum":["open","search","youtube","naver","github","translate"]},"query":{"type":"string"},"url":{"type":"string"}},"required":["action"]}}},
    {"type":"function","function":{"name":"window_manager","description":"창 관리: 최소화/최대화/닫기/앞으로","parameters":{"type":"object","properties":{"action":{"type":"string","enum":["minimize_all","maximize","focus","list_windows","close_window"]},"window_title":{"type":"string"}},"required":["action"]}}},
    {"type":"function","function":{"name":"file_search","description":"파일 검색 (이름/확장자/내용으로 찾기)","parameters":{"type":"object","properties":{"query":{"type":"string"},"search_path":{"type":"string"},"ext":{"type":"string"},"search_content":{"type":"boolean"}},"required":["query"]}}},
    {"type":"function","function":{"name":"download_file","description":"URL에서 파일 다운로드","parameters":{"type":"object","properties":{"url":{"type":"string"},"save_path":{"type":"string"}},"required":["url"]}}},
    {"type":"function","function":{"name":"compress_files","description":"파일 압축/해제 (zip)","parameters":{"type":"object","properties":{"action":{"type":"string","enum":["compress","extract"]},"path":{"type":"string"},"output":{"type":"string"}},"required":["action","path"]}}},
    {"type":"function","function":{"name":"notify","description":"Windows 알림 팝업 띄우기","parameters":{"type":"object","properties":{"title":{"type":"string"},"message":{"type":"string"}},"required":["title","message"]}}},
]

# ──────────────────────── 도구 실행기 ────────────────────────
def execute_tool(name, args, st=None):
    def status(m):
        if st: st(m)
    try:
        if name == "read_file":
            c = Path(args["path"]).read_text(encoding="utf-8", errors="replace")
            return c[:8000] + ("...[잘림]" if len(c)>8000 else "")
        elif name == "write_file":
            p = Path(args["path"]); p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(args["content"], encoding="utf-8"); return f"✅ 저장: {args['path']}"
        elif name == "append_file":
            with open(args["path"],"a",encoding="utf-8") as f: f.write(args["content"])
            return f"✅ 추가: {args['path']}"
        elif name == "list_directory":
            base=args["path"]; max_d=args.get("depth",2); lines=[]
            for root,dirs,files in os.walk(base):
                d=root.replace(base,"").count(os.sep)
                if d>max_d: dirs.clear(); continue
                lines.append("│   "*d+f"├── {os.path.basename(root)}/")
                for f in sorted(files):
                    sz=os.path.getsize(os.path.join(root,f))
                    lines.append("│   "*(d+1)+f"├── {f} [{sz//1024}KB]")
            return "\n".join(lines) or "빈 디렉토리"
        elif name == "copy_file":
            shutil.copy2(args["src"],args["dst"]); return f"✅ 복사 완료"
        elif name == "delete_file":
            p=Path(args["path"])
            (p.unlink() if p.is_file() else p.rmdir()); return f"✅ 삭제"
        elif name == "run_command":
            status(f"명령 실행: {args['command'][:60]}")
            return run_cmd(args["command"], args.get("shell","cmd"), args.get("timeout",120))
        elif name == "edit_pptx":
            status(f"PPTX 편집: {args.get('operation','')}")
            return edit_pptx_tool(args)
        elif name == "read_pdf":
            if not LIBS["pdf"]: return "❌ pip install pdfplumber"
            with pdfplumber.open(args["path"]) as pdf:
                return "\n\n".join(f"[p{i+1}]\n{p.extract_text()or''}" for i,p in enumerate(pdf.pages))[:6000]
        elif name == "read_word":
            if not LIBS["docx"]: return "❌ pip install python-docx"
            return "\n".join(p.text for p in DocxDocument(args["path"]).paragraphs)[:6000]
        elif name == "write_word":
            if not LIBS["docx"]: return "❌ pip install python-docx"
            doc=DocxDocument()
            for line in args["content"].split("\n"): doc.add_paragraph(line)
            doc.save(args["path"]); return f"✅ Word 저장"
        elif name == "read_excel":
            if not LIBS["xlsx"]: return "❌ pip install openpyxl"
            wb=openpyxl.load_workbook(args["path"])
            out=[]
            for sn in wb.sheetnames:
                out.append(f"[{sn}]")
                for row in wb[sn].iter_rows(values_only=True):
                    out.append("\t".join(str(c) if c else "" for c in row))
            return "\n".join(out)[:5000]
        elif name == "analyze_image":
            if not LIBS["PIL"]: return "❌ pip install pillow"
            status("이미지 압축 중...")
            b64,mime=compress_img(args["path"])
            status("Vision AI 분석 중...")
            return f"[Vision AI]\n{call_vision(b64,mime,args.get('question','이 이미지를 설명해줘'))}"
        elif name == "see_screen":
            if not LIBS["PIL"]: return "❌ pip install pillow"
            screen=ImageGrab.grab()
            sp=args.get("save_path",f"screen_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png")
            screen.save(sp)
            b64,mime=compress_img(screen)
            return f"[화면 {screen.width}x{screen.height}]\n{call_vision(b64,mime,args.get('question','화면 설명'))}"
        elif name == "capture_webcam":
            if not LIBS["cv2"]: return "❌ pip install opencv-python"
            cap=cv2.VideoCapture(0)
            if not cap.isOpened(): return "❌ 웹캠 없음"
            ret,frame=cap.read(); cap.release()
            if not ret: return "❌ 캡처 실패"
            b64,mime=compress_img(Image.fromarray(cv2.cvtColor(frame,cv2.COLOR_BGR2RGB)))
            return f"[웹캠]\n{call_vision(b64,mime,args.get('question','설명해줘'))}"
        elif name == "analyze_video":
            if not LIBS["cv2"]: return "❌ pip install opencv-python"
            cap=cv2.VideoCapture(args["path"])
            if not cap.isOpened(): return "❌ 비디오 열기 실패"
            total=int(cap.get(cv2.CAP_PROP_FRAME_COUNT)); fps=cap.get(cv2.CAP_PROP_FPS) or 30
            num=args.get("num_frames",5)
            results=[f"[{Path(args['path']).name} | {total/fps:.1f}초]"]
            for i,fi in enumerate([int(j*total/num) for j in range(num)]):
                cap.set(cv2.CAP_PROP_POS_FRAMES,fi)
                ret,frame=cap.read()
                if not ret: continue
                t=fi/fps
                b64,mime=compress_img(Image.fromarray(cv2.cvtColor(frame,cv2.COLOR_BGR2RGB)),max_side=960)
                status(f"프레임 {i+1}/{num} ({t:.1f}초)")
                q = f'{t:.1f}초: {args["question"]}'
                results.append(f"\n[{t:.1f}초] {call_vision(b64,mime,q)}")
            cap.release(); return "\n".join(results)
        elif name == "analyze_audio":
            path=args["path"]
            if args.get("task")=="transcribe":
                if not LIBS["speech"]: return "❌ pip install SpeechRecognition"
                recognizer=sr.Recognizer()
                wav=path if path.lower().endswith(".wav") else path.rsplit(".",1)[0]+"_conv.wav"
                if not path.lower().endswith(".wav"):
                    subprocess.run(["ffmpeg","-i",path,wav,"-y"],capture_output=True,timeout=30)
                with sr.AudioFile(wav) as src: audio=recognizer.record(src)
                try: return f"[전사]\n{recognizer.recognize_google(audio,language='ko-KR')}"
                except: return "❌ 인식 실패"
            info={"파일":path,"크기":f"{os.path.getsize(path):,}B","확장자":Path(path).suffix}
            try:
                from mutagen import File as MF; mf=MF(path)
                if mf and mf.info: info["길이"]=f"{mf.info.length:.1f}초"
            except: pass
            return json.dumps(info,ensure_ascii=False,indent=2)
        elif name == "edit_image":
            if not LIBS["PIL"]: return "❌ pip install pillow"
            img=Image.open(args["path"]); op=args["operation"]; out=args.get("output_path",args["path"])
            ops={"grayscale":lambda i:i.convert("L"),"blur":lambda i:i.filter(ImageFilter.GaussianBlur(args.get("factor",2))),"sharpen":lambda i:i.filter(ImageFilter.SHARPEN),"flip":lambda i:i.transpose(Image.FLIP_LEFT_RIGHT),"rotate":lambda i:i.rotate(args.get("angle",90),expand=True),"brightness":lambda i:ImageEnhance.Brightness(i).enhance(args.get("factor",1.5)),"contrast":lambda i:ImageEnhance.Contrast(i).enhance(args.get("factor",1.5))}
            if op in ops: img=ops[op](img)
            elif op=="resize": img=img.resize((args.get("width",img.width),args.get("height",img.height)),Image.LANCZOS)
            elif op=="add_text": ImageDraw.Draw(img).text((10,10),args.get("text",""),fill=(255,255,0))
            elif op=="crop": img=img.crop((args.get("x1",0),args.get("y1",0),args.get("x2",100),args.get("y2",100)))
            img.save(out); return f"✅ 이미지 편집 [{op}]: {out}"
        elif name == "open_browser":
            webbrowser.open(args["url"]); return f"✅ 브라우저: {args['url']}"
        elif name == "google_search":
            url=f"https://www.google.com/search?q={quote(args['query'])}"
            if args.get("open_browser",True): webbrowser.open(url); return f"✅ Google: {args['query']}"
            r=requests.get(url,headers={"User-Agent":"Mozilla/5.0"},timeout=10)
            return re.sub(r'\s+',' ',re.sub(r'<[^>]+>',' ',re.sub(r'<script.*?</script>','',r.text,flags=re.DOTALL))).strip()[:3000]
        elif name == "fetch_url":
            r=requests.get(args["url"],headers={"User-Agent":"Mozilla/5.0"},timeout=15)
            t=re.sub(r'<style.*?</style>','',re.sub(r'<script.*?</script>','',r.text,flags=re.DOTALL),flags=re.DOTALL)
            return re.sub(r'\s+',' ',re.sub(r'<[^>]+>',' ',t)).strip()[:args.get("max_chars",3000)]
        elif name == "send_email":
            if not EMAIL_CONFIG["email"]: return "❌ 이메일 설정 필요"
            msg=MIMEMultipart(); msg["From"]=EMAIL_CONFIG["email"]; msg["To"]=args["to"]; msg["Subject"]=args["subject"]
            msg.attach(MIMEText(args["body"],"plain","utf-8"))
            with smtplib.SMTP(EMAIL_CONFIG["smtp_host"],EMAIL_CONFIG["smtp_port"]) as s:
                s.ehlo(); s.starttls(); s.login(EMAIL_CONFIG["email"],EMAIL_CONFIG["password"]); s.send_message(msg)
            return f"✅ 이메일 발송 → {args['to']}"
        elif name == "open_email_client":
            c=args.get("client","gmail"); to=args.get("to",""); sub=quote(args.get("subject","")); body=quote(args.get("body",""))
            urls={"gmail":f"https://mail.google.com/mail/?view=cm&to={to}&su={sub}&body={body}","outlook":f"https://outlook.live.com/mail/0/deeplink/compose?to={to}&subject={sub}&body={body}"}
            webbrowser.open(urls.get(c,f"mailto:{to}?subject={sub}&body={body}")); return "✅ 이메일 클라이언트 열기"
        elif name == "launch_app":
            app=args["app"].lower(); extra=args.get("args","")
            APP_MAP={
                # 브라우저
                "chrome":["chrome","C:/Program Files/Google/Chrome/Application/chrome.exe","C:/Program Files (x86)/Google/Chrome/Application/chrome.exe"],
                "edge":["msedge","microsoft-edge"],
                "firefox":["firefox","C:/Program Files/Mozilla Firefox/firefox.exe"],
                "whale":["whale","C:/Program Files/Naver/Naver Whale/whale.exe"],
                "opera":["opera"],
                # 오피스
                "excel":["excel","EXCEL.EXE"],"word":["winword","WINWORD.EXE"],
                "powerpoint":["powerpnt","POWERPNT.EXE"],"pptx":["powerpnt"],
                "outlook":["outlook"],"onenote":["onenote"],
                "hwp":["hwp","C:/Program Files/HNC/Office NEO/HOffice110/Bin/Hwp.exe"],
                "hancom":["hwp"],
                # 개발도구
                "vscode":["code","C:/Users/"+os.environ.get("USERNAME","user")+"/AppData/Local/Programs/Microsoft VS Code/Code.exe"],
                "pycharm":["pycharm","C:/Program Files/JetBrains/PyCharm/bin/pycharm64.exe"],
                "notepad":["notepad"],"notepad++":["notepad++","C:/Program Files/Notepad++/notepad++.exe"],
                "git":["git-bash","C:/Program Files/Git/git-bash.exe"],
                "cmd":["cmd"],"powershell":["powershell"],"terminal":["wt","powershell"],
                # 미디어
                "vlc":["vlc","C:/Program Files/VideoLAN/VLC/vlc.exe"],
                "spotify":["spotify","C:/Users/"+os.environ.get("USERNAME","user")+"/AppData/Roaming/Spotify/Spotify.exe"],
                "itunes":["itunes"],"windows media player":["wmplayer"],
                # 커뮤니케이션
                "discord":["discord","C:/Users/"+os.environ.get("USERNAME","user")+"/AppData/Local/Discord/app-*/Discord.exe"],
                "slack":["slack","C:/Users/"+os.environ.get("USERNAME","user")+"/AppData/Local/slack/slack.exe"],
                "teams":["teams","msteams"],
                "kakaotalk":["kakaotalk","C:/Program Files (x86)/Kakao/KakaoTalk/KakaoTalk.exe"],
                "zoom":["zoom","C:/Users/"+os.environ.get("USERNAME","user")+"/AppData/Roaming/Zoom/bin/Zoom.exe"],
                "telegram":["telegram","C:/Users/"+os.environ.get("USERNAME","user")+"/AppData/Roaming/Telegram Desktop/Telegram.exe"],
                # 창작/디자인
                "photoshop":["photoshop","C:/Program Files/Adobe/Adobe Photoshop/Photoshop.exe"],
                "illustrator":["illustrator"],"premiere":["premiere pro"],"after effects":["afterfx"],
                "blender":["blender","C:/Program Files/Blender Foundation/Blender/blender.exe"],
                "obs":["obs64","C:/Program Files/obs-studio/bin/64bit/obs64.exe"],
                "figma":["figma"],
                # 게임/스토어
                "steam":["steam","C:/Program Files (x86)/Steam/steam.exe"],
                "epic":["epicgameslauncher","C:/Program Files (x86)/Epic Games/Launcher/Portal/Binaries/Win32/EpicGamesLauncher.exe"],
                # 유틸리티
                "calc":["calc"],"calculator":["calc"],
                "paint":["mspaint"],"paint3d":["mspaint3d"],
                "snipping":["snippingtool","SnippingTool.exe"],
                "taskmgr":["taskmgr"],"taskmanager":["taskmgr"],
                "explorer":["explorer"],"file explorer":["explorer"],
                "control":["control"],"settings":["ms-settings:"],
                "winrar":["winrar","C:/Program Files/WinRAR/WinRAR.exe"],
                "7zip":["7zfm","C:/Program Files/7-Zip/7zFM.exe"],
                "filezilla":["filezilla","C:/Program Files/FileZilla FTP Client/filezilla.exe"],
                "putty":["putty"],"winscp":["winscp"],
                "python":["python","python3"],"jupyter":["jupyter","jupyter-notebook"],
            }
            candidates = APP_MAP.get(app, [app])
            for c in candidates:
                try:
                    if "*" in c:  # glob 패턴
                        matches = glob.glob(c)
                        if matches: c = matches[-1]
                        else: continue
                    cmd = f'"{c}" {extra}' if extra else c
                    subprocess.Popen(cmd, shell=True)
                    return f"✅ 앱 실행: {c}"
                except: pass
            # 마지막 수단: Windows Start 명령
            try:
                subprocess.Popen(f"start {app}", shell=True)
                return f"✅ Start 명령으로 실행: {app}"
            except: pass
            return f"❌ 앱을 찾을 수 없음: {app} (직접 경로를 알려주세요)"
        elif name == "list_running_apps":
            if not LIBS["psutil"]: return subprocess.run("tasklist",capture_output=True,text=True,shell=True).stdout[:3000]
            filt=args.get("filter","").lower()
            procs=[]
            for p in psutil.process_iter(['pid','name','memory_info']):
                try:
                    n=p.info['name'] or ""
                    if filt and filt not in n.lower(): continue
                    procs.append(f"PID:{p.info['pid']:6d}  {p.info['memory_info'].rss//1024//1024:5d}MB  {n}")
                except: pass
            return "\n".join(procs[:60])
        elif name == "kill_process":
            t=args["name_or_pid"]
            subprocess.run(f"taskkill /{'PID' if t.isdigit() else 'IM'} {t} /F",shell=True); return f"✅ 종료: {t}"
        elif name == "mouse_control":
            if not LIBS["pyautogui"]: return "❌ pip install pyautogui"
            a=args["action"]; x,y=args.get("x",0),args.get("y",0)
            {"move":lambda:pyautogui.moveTo(x,y,duration=args.get("duration",0.3)),"click":lambda:pyautogui.click(x,y),"double_click":lambda:pyautogui.doubleClick(x,y),"right_click":lambda:pyautogui.rightClick(x,y),"drag":lambda:pyautogui.dragTo(args.get("x2",x),args.get("y2",y),duration=0.3),"scroll":lambda:pyautogui.scroll(args.get("dy",3),x=x,y=y)}.get(a,lambda:None)()
            return f"✅ 마우스 {a}"
        elif name == "keyboard_control":
            if not LIBS["pyautogui"]: return "❌ pip install pyautogui"
            a=args["action"]
            if a=="type": pyautogui.write(args.get("text",""),interval=0.03)
            elif a=="hotkey": pyautogui.hotkey(*[k.strip() for k in args.get("key","").split("+")])
            elif a=="press": pyautogui.press(args.get("key","enter"))
            return f"✅ 키보드 {a}"
        elif name == "take_screenshot":
            if not LIBS["PIL"]: return "❌ pip install pillow"
            p=args.get("save_path",f"shot_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png")
            screen=ImageGrab.grab(); screen.save(p); return f"✅ 스크린샷: {p}"
        elif name == "get_system_info":
            import platform
            info={"OS":platform.system(),"Version":platform.version(),"Python":sys.version.split()[0],"CWD":os.getcwd()}
            if LIBS["psutil"]:
                info["CPU"]=f"{psutil.cpu_percent(1)}%"
                m=psutil.virtual_memory()
                info["RAM"]=f"{m.used//1024//1024:,}MB/{m.total//1024//1024:,}MB ({m.percent}%)"
            return json.dumps(info,ensure_ascii=False,indent=2)
        elif name == "clipboard":
            if args["action"]=="read":
                out, _, _ = run_ps("[Console]::OutputEncoding = [System.Text.Encoding]::UTF8\nGet-Clipboard")
                return out or "(비어있음)"
            ps_set = f'[Console]::OutputEncoding = [System.Text.Encoding]::UTF8\nSet-Clipboard -Value @''\n{args.get("text","")}\n''@'
            run_ps(ps_set)
            return "✅ 클립보드 복사"

        elif name == "create_python":
            path = args["path"]
            code = args["code"]
            Path(path).parent.mkdir(parents=True, exist_ok=True)
            Path(path).write_text(code, encoding="utf-8")
            result = f"✅ Python 파일 생성: {path}\n코드 길이: {len(code.splitlines())}줄"
            if args.get("run", False):
                extra_args = args.get("args", "")
                try:
                    out = subprocess.run(
                        f'python "{path}" {extra_args}',
                        shell=True, capture_output=True, text=True,
                        timeout=args.get("timeout", 60), encoding="utf-8", errors="replace"
                    )
                    result += f"\n\n[실행 결과]\n{out.stdout or ''}{out.stderr or ''}"
                except subprocess.TimeoutExpired:
                    result += "\n⚠️ 실행 시간 초과"
                except Exception as e:
                    result += f"\n❌ 실행 오류: {e}"
            return result

        elif name == "run_python":
            code = args["code"]
            try:
                # utf-8 출력 강제 헤더 삽입
                header = "# -*- coding: utf-8 -*-\nimport sys, io\nsys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')\nsys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')\n"
                with tempfile.NamedTemporaryFile(mode="w", suffix=".py", delete=False, encoding="utf-8") as f:
                    f.write(header + code); tmp = f.name
                out = subprocess.run(
                    f'python -X utf8 "{tmp}"', shell=True, capture_output=True,
                    timeout=args.get("timeout", 30),
                )
                os.unlink(tmp)
                stdout = out.stdout.decode("utf-8", errors="replace") if out.stdout else ""
                stderr = out.stderr.decode("utf-8", errors="replace") if out.stderr else ""
                return f"[출력]\n{stdout}\n[오류]\n{stderr}" if stderr.strip() else f"[출력]\n{stdout or '(없음)'}"
            except subprocess.TimeoutExpired:
                return "⚠️ 타임아웃"
            except Exception as e:
                return f"❌ 실행 오류: {e}"

        elif name == "open_file_dialog":
            path = args["path"]
            action = args.get("action", "open_folder")
            if action == "open_folder":
                subprocess.Popen(f'explorer "{path}"', shell=True)
                return f"✅ 폴더 열기: {path}"
            elif action == "open_file":
                os.startfile(path)
                return f"✅ 파일 열기: {path}"
            elif action == "reveal":
                subprocess.Popen(f'explorer /select,"{path}"', shell=True)
                return f"✅ 파일 위치 표시: {path}"

        elif name == "create_pptx":
            if not LIBS["pptx"]: return "❌ pip install python-pptx"
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            prs = Presentation()
            theme = args.get("theme", "dark")
            THEMES = {
                "dark":  {"bg": (7,9,15), "title": (0,212,255), "body": (238,244,255)},
                "blue":  {"bg": (13,27,62), "title": (100,180,255), "body": (220,235,255)},
                "white": {"bg": (255,255,255), "title": (30,60,120), "body": (50,50,50)},
                "green": {"bg": (5,20,10), "title": (0,255,136), "body": (200,255,220)},
            }
            tc = THEMES.get(theme, THEMES["dark"])
            from pptx.util import Emu
            from pptx.oxml.ns import qn
            import lxml.etree as etree
            for slide_data in args.get("slides", []):
                layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(layout)
                # 배경색
                bg = slide.background; fill = bg.fill; fill.solid()
                fill.fore_color.rgb = RGBColor(*tc["bg"])
                # 제목
                if slide.shapes.title:
                    slide.shapes.title.text = slide_data.get("title", "")
                    for para in slide.shapes.title.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.color.rgb = RGBColor(*tc["title"])
                            run.font.size = Pt(32)
                            run.font.bold = True
                # 내용
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 1:
                        ph.text = slide_data.get("content", "")
                        for para in ph.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.color.rgb = RGBColor(*tc["body"])
                                run.font.size = Pt(18)
            path = args["path"]
            Path(path).parent.mkdir(parents=True, exist_ok=True)
            prs.save(path)
            return f"✅ PPTX 생성: {path} ({len(args.get('slides',[]))}슬라이드, {theme} 테마)"

        elif name == "browser_control":
            action = args.get("action", "open")
            query = args.get("query", "")
            url = args.get("url", "")
            urls = {
                "search":    f"https://www.google.com/search?q={quote(query)}",
                "youtube":   f"https://www.youtube.com/results?search_query={quote(query)}",
                "naver":     f"https://search.naver.com/search.naver?query={quote(query)}",
                "github":    f"https://github.com/search?q={quote(query)}&type=repositories",
                "translate": f"https://translate.google.com/?text={quote(query)}&sl=ko&tl=en",
                "open":      url or "https://www.google.com",
            }
            target = urls.get(action, url)
            webbrowser.open(target)
            return f"✅ 브라우저 {action}: {target}"

        elif name == "window_manager":
            action = args.get("action")
            title = args.get("window_title", "")
            if action == "minimize_all":
                subprocess.run("powershell (New-Object -ComObject Shell.Application).MinimizeAll()", shell=True)
                return "✅ 모든 창 최소화"
            elif action == "list_windows":
                out, err, _ = run_ps(
                    "[Console]::OutputEncoding = [System.Text.Encoding]::UTF8\n"
                    "Get-Process | Where-Object {$_.MainWindowTitle -ne ''} | "
                    "Select-Object Name,Id,MainWindowTitle | Format-Table -AutoSize"
                )
                return out or "실행 중인 창 없음"
            elif action in ("maximize", "focus", "close_window") and title:
                ps = f"""
Add-Type @'
using System;
using System.Runtime.InteropServices;
public class Win {{
    [DllImport("user32.dll")] public static extern bool ShowWindow(IntPtr hWnd, int nCmdShow);
    [DllImport("user32.dll")] public static extern bool SetForegroundWindow(IntPtr hWnd);
    [DllImport("user32.dll")] public static extern bool PostMessage(IntPtr hWnd, uint Msg, int wParam, int lParam);
}}
'@ -Language CSharp
$p = Get-Process | Where-Object {{$_.MainWindowTitle -like "*{title}*"}} | Select-Object -First 1
if ($p) {{
    if ("{action}" -eq "maximize") {{ [Win]::ShowWindow($p.MainWindowHandle, 3) }}
    elseif ("{action}" -eq "focus") {{ [Win]::SetForegroundWindow($p.MainWindowHandle) }}
    elseif ("{action}" -eq "close_window") {{ $p.CloseMainWindow() }}
    "✅ {action}: $($p.MainWindowTitle)"
}} else {{ "❌ 창 없음: {title}" }}
"""
                out, err, rc = run_ps(ps)
                return out or err
            return f"✅ {action} 완료"

        elif name == "file_search":
            query = args["query"]; base = args.get("search_path", os.path.expanduser("~"))
            ext = args.get("ext", ""); results = []
            pattern = f"**/*{query}*{ext}" if ext else f"**/*{query}*"
            try:
                for p in Path(base).glob(pattern):
                    results.append(str(p))
                    if len(results) >= 50: break
            except: pass
            if not results:
                # fallback: where 명령
                out = subprocess.run(f'where /R "{base}" *{query}*', capture_output=True, text=True, shell=True, encoding="utf-8", errors="replace")
                results = out.stdout.strip().split("\n")[:30]
            return "\n".join(results) or f"❌ 검색 결과 없음: {query}"

        elif name == "download_file":
            url = args["url"]; save_path = args.get("save_path", url.split("/")[-1])
            status(f"다운로드 중: {url[:60]}")
            r = requests.get(url, stream=True, timeout=60, headers={"User-Agent":"Mozilla/5.0"})
            Path(save_path).parent.mkdir(parents=True, exist_ok=True)
            with open(save_path, "wb") as f:
                for chunk in r.iter_content(8192): f.write(chunk)
            return f"✅ 다운로드 완료: {save_path} ({os.path.getsize(save_path):,}B)"

        elif name == "compress_files":
            action = args["action"]; path = args["path"]; output = args.get("output", "")
            if action == "compress":
                out_zip = output or path + ".zip"
                subprocess.run(f'powershell Compress-Archive -Path "{path}" -DestinationPath "{out_zip}" -Force', shell=True)
                return f"✅ 압축: {out_zip}"
            elif action == "extract":
                out_dir = output or Path(path).parent / Path(path).stem
                subprocess.run(f'powershell Expand-Archive -Path "{path}" -DestinationPath "{out_dir}" -Force', shell=True)
                return f"✅ 압축 해제: {out_dir}"

        elif name == "notify":
            title = args.get("title","JihunAI"); msg = args.get("message","")
            ps = f'[Windows.UI.Notifications.ToastNotificationManager, Windows.UI.Notifications, ContentType = WindowsRuntime] | Out-Null\n$t = [Windows.UI.Notifications.ToastNotificationManager]::GetTemplateContent([Windows.UI.Notifications.ToastTemplateType]::ToastText02)\n$t.GetElementsByTagName("text")[0].AppendChild($t.CreateTextNode("{title}")) | Out-Null\n$t.GetElementsByTagName("text")[1].AppendChild($t.CreateTextNode("{msg}")) | Out-Null\n[Windows.UI.Notifications.ToastNotificationManager]::CreateToastNotifier("JihunAI").Show([Windows.UI.Notifications.ToastNotification]::new($t))'
            run_ps(ps)
            return f"✅ 알림: {title}"

        return f"❌ 알 수 없는 도구: {name}"
    except Exception as e:
        import traceback
        return f"❌ [{name}] {e}\n{traceback.format_exc()[-400:]}"

# ═══════════════════════════════════════════════════════════════
#  JihunAI v1.1 앱
# ═══════════════════════════════════════════════════════════════
class JihunAI:
    def __init__(self, root):
        self.root = root
        self.root.title("JihunAI v1.4")
        self.root.geometry("1540x980")
        self.root.configure(bg=C["bg"])
        self.root.minsize(1100, 720)

        self.conversation    = [{"role":"system","content":SYSTEM_PROMPT}]
        self.current_chat_id = None
        self.is_processing   = False
        self._tick           = 0
        self._tasks          = []
        self._current_task   = None
        self._active_model   = DEFAULT_MODEL
        self._fallback_idx   = 0  # fallback 체인 인덱스
        self._retry_counts   = {}  # task_id → retry_count
        self._bot_stats      = {m: {"calls":0,"errors":0,"success":0} for m in BOT_ROSTER}

        self._build_ui()
        self.load_chat_list()
        self.load_latest_chat()
        self.root.protocol("WM_DELETE_WINDOW", self._on_close)
        self._animation_loop()

    # ─────────────── UI ───────────────
    def _build_ui(self):
        self.root.grid_columnconfigure(1, weight=1)
        self.root.grid_rowconfigure(0, weight=1)
        self._build_sidebar()
        self._build_center()

    def _build_sidebar(self):
        sb = tk.Frame(self.root, bg=C["panel"], width=270)
        sb.grid(row=0, column=0, sticky="nsew")
        sb.grid_propagate(False)
        sb.grid_rowconfigure(5, weight=1)
        sb.grid_columnconfigure(0, weight=1)

        # 로고
        logo = tk.Frame(sb, bg=C["panel"], pady=14)
        logo.grid(row=0, column=0, sticky="ew")
        self._orb = tk.Canvas(logo, width=40, height=40, bg=C["panel"], highlightthickness=0)
        self._orb.pack(side="left", padx=(14,8))
        lbl = tk.Frame(logo, bg=C["panel"]); lbl.pack(side="left")
        tk.Label(lbl, text="JihunAI", font=("Consolas",16,"bold"), fg=C["cyan"], bg=C["panel"]).pack(anchor="w")
        tk.Label(lbl, text="v1.4 · GPT우선 · Vision수정 · 완전체", font=("Consolas",7), fg=C["dim"], bg=C["panel"]).pack(anchor="w")

        tk.Frame(sb, bg=C["border2"], height=1).grid(row=1, column=0, sticky="ew", padx=10)

        btns = tk.Frame(sb, bg=C["panel"])
        btns.grid(row=2, column=0, sticky="ew", padx=10, pady=6)
        for txt, cmd, col, side in [("＋ NEW",self._new_chat,C["cyan"],"left"),
                                      ("✕ DEL",self._del_chat,C["red"],"left"),
                                      ("⚙",self._config,C["gold"],"right")]:
            tk.Button(btns, text=txt, command=cmd, bg=C["bg3"], fg=col,
                      font=("Consolas",8,"bold"), relief="flat", bd=0, padx=9, pady=5,
                      activebackground=C["border2"], activeforeground=col, cursor="hand2").pack(side=side, padx=2)

        # 멀티봇 현황
        bf = tk.Frame(sb, bg=C["bg2"], padx=10, pady=8)
        bf.grid(row=3, column=0, sticky="ew", padx=10, pady=4)
        tk.Label(bf, text="◈ MULTI-BOT STATUS", font=("Consolas",7,"bold"), fg=C["purple"], bg=C["bg2"]).pack(anchor="w")
        self._bot_labels = {}
        for m_key in list(BOT_ROSTER.keys())[:6]:
            _, _, desc = BOT_ROSTER[m_key]
            r = tk.Frame(bf, bg=C["bg2"]); r.pack(fill="x", pady=1)
            dot = tk.Label(r, text="●", fg=C["dim"], bg=C["bg2"], font=("Consolas",7)); dot.pack(side="left")
            tk.Label(r, text=f"  {desc[:28]}", fg=C["dim"], bg=C["bg2"], font=("Consolas",7)).pack(side="left")
            self._bot_labels[m_key] = dot

        # 태스크 패널
        tp = tk.Frame(sb, bg=C["bg2"], padx=10, pady=8)
        tp.grid(row=4, column=0, sticky="ew", padx=10, pady=4)
        tk.Label(tp, text="◈ TASK PROGRESS", font=("Consolas",7,"bold"), fg=C["cyan"], bg=C["bg2"]).pack(anchor="w")
        self._task_frame = tk.Frame(tp, bg=C["bg2"])
        self._task_frame.pack(fill="x", pady=(4,0))

        # 채팅 목록
        lf = tk.Frame(sb, bg=C["panel"])
        lf.grid(row=5, column=0, sticky="nsew", padx=10, pady=(4,0))
        lf.grid_rowconfigure(1, weight=1); lf.grid_columnconfigure(0, weight=1)
        tk.Label(lf, text="MISSION LOG", font=("Consolas",7), fg=C["dim"], bg=C["panel"], anchor="w").grid(row=0, column=0, sticky="ew")
        self._chat_lb = tk.Listbox(lf, bg=C["bg2"], fg=C["mid"],
                                    selectbackground=C["border2"], selectforeground=C["cyan"],
                                    font=("Consolas",8), borderwidth=0, highlightthickness=1,
                                    highlightbackground=C["border"], highlightcolor=C["cyan"],
                                    activestyle="none", relief="flat")
        self._chat_lb.grid(row=1, column=0, sticky="nsew")
        tk.Scrollbar(lf, orient="vertical", command=self._chat_lb.yview).grid(row=1, column=1, sticky="ns")
        self._chat_lb.bind("<<ListboxSelect>>", self._on_chat_sel)

        # 모듈 상태
        mf = tk.Frame(sb, bg=C["bg2"], padx=10, pady=6)
        mf.grid(row=6, column=0, sticky="ew", padx=10, pady=(4,10))
        tk.Label(mf, text="MODULE", font=("Consolas",7), fg=C["dim"], bg=C["bg2"]).pack(anchor="w")
        for nm,key,col in [("Vision","PIL",C["cyan"]),("PPTX","pptx",C["gold"]),("PyAutoGUI","pyautogui",C["green"]),("OpenCV","cv2",C["purple"]),("Psutil","psutil",C["orange"])]:
            r=tk.Frame(mf,bg=C["bg2"]); r.pack(fill="x",pady=1)
            tk.Label(r,text="●",fg=col if LIBS.get(key) else C["red"],bg=C["bg2"],font=("Consolas",7)).pack(side="left")
            tk.Label(r,text=f"  {nm:<12}{'ON' if LIBS.get(key) else 'OFF'}",fg=C["mid"] if LIBS.get(key) else C["dim"],bg=C["bg2"],font=("Consolas",7)).pack(side="left")

    def _build_center(self):
        center = tk.Frame(self.root, bg=C["bg"])
        center.grid(row=0, column=1, sticky="nsew")
        center.grid_rowconfigure(1, weight=1); center.grid_columnconfigure(0, weight=1)

        hdr = tk.Frame(center, bg=C["panel"], height=56)
        hdr.grid(row=0, column=0, sticky="ew"); hdr.grid_propagate(False)
        hdr.grid_columnconfigure(1, weight=1)
        left = tk.Frame(hdr, bg=C["panel"]); left.grid(row=0, column=0, padx=20, pady=8, sticky="w")
        tk.Label(left, text="◈ JihunAI", font=("Consolas",13,"bold"), fg=C["cyan"], bg=C["panel"]).pack(side="left")
        self._thinking_lbl = tk.Label(left, text="", font=("Consolas",9), fg=C["thinking"], bg=C["panel"])
        self._thinking_lbl.pack(side="left", padx=(16,0))
        # 현재 봇 표시
        self._bot_lbl = tk.Label(left, text="", font=("Consolas",8), fg=C["purple"], bg=C["panel"])
        self._bot_lbl.pack(side="left", padx=(10,0))

        right = tk.Frame(hdr, bg=C["panel"]); right.grid(row=0, column=2, padx=20, sticky="e")
        tk.Label(right, text="MAIN BOT", font=("Consolas",7), fg=C["dim"], bg=C["panel"]).pack(anchor="e")
        self._model_var = tk.StringVar(value=DEFAULT_MODEL)  # gpt-5.4-pro
        self._style_combo()
        ttk.Combobox(right, textvariable=self._model_var, values=list(BOT_ROSTER.keys()),
                     state="readonly", width=22, font=("Consolas",9)).pack()

        chat_wrap = tk.Frame(center, bg=C["bg"])
        chat_wrap.grid(row=1, column=0, sticky="nsew", padx=10, pady=(6,0))
        chat_wrap.grid_rowconfigure(0, weight=1); chat_wrap.grid_columnconfigure(0, weight=1)
        self._chat = scrolledtext.ScrolledText(
            chat_wrap, wrap=tk.WORD, state="disabled", font=("Consolas",11),
            bg=C["bg"], fg=C["white"], insertbackground=C["cyan"],
            selectbackground=C["bg3"], relief="flat", borderwidth=0, padx=20, pady=16, spacing1=3, spacing3=3)
        self._chat.grid(row=0, column=0, sticky="nsew")
        self._setup_tags()
        self._build_input(center)

    def _setup_tags(self):
        t = self._chat
        t.tag_configure("ts",       foreground=C["dim"],     font=("Consolas",7))
        t.tag_configure("u_lbl",    foreground=C["cyan"],    font=("Consolas",8,"bold"))
        t.tag_configure("u_txt",    foreground=C["white"],   font=("Consolas",11), lmargin1=24, lmargin2=24)
        t.tag_configure("a_lbl",    foreground=C["green"],   font=("Consolas",8,"bold"))
        t.tag_configure("a_txt",    foreground=C["white"],   font=("Consolas",11), lmargin1=24, lmargin2=24)
        t.tag_configure("tc_lbl",   foreground=C["gold"],    font=("Consolas",8,"bold"))
        t.tag_configure("tc_txt",   foreground=C["gold"],    font=("Consolas",9,"italic"), lmargin1=24, lmargin2=24)
        t.tag_configure("tr_txt",   foreground=C["mid"],     font=("Consolas",9), lmargin1=24, lmargin2=24)
        t.tag_configure("err_txt",  foreground=C["red"],     font=("Consolas",10), lmargin1=24, lmargin2=24)
        t.tag_configure("sys_txt",  foreground=C["dim"],     font=("Consolas",9,"italic"), justify="center")
        t.tag_configure("plan_lbl", foreground=C["purple"],  font=("Consolas",8,"bold"))
        t.tag_configure("plan_txt", foreground="#cc88ff",    font=("Consolas",10), lmargin1=24, lmargin2=24)
        t.tag_configure("task_hdr", foreground=C["cyan"],    font=("Consolas",9,"bold"), lmargin1=24, lmargin2=24)
        t.tag_configure("task_ok",  foreground=C["green"],   font=("Consolas",9), lmargin1=32, lmargin2=32)
        t.tag_configure("task_fail",foreground=C["red"],     font=("Consolas",9), lmargin1=32, lmargin2=32)
        t.tag_configure("retry_lbl",foreground=C["retry"],   font=("Consolas",8,"bold"))
        t.tag_configure("fb_lbl",   foreground=C["fallback"],font=("Consolas",8,"bold"))

    def _build_input(self, parent):
        wrap = tk.Frame(parent, bg=C["panel"])
        wrap.grid(row=2, column=0, sticky="ew", padx=10, pady=(4,10))
        wrap.grid_columnconfigure(0, weight=1)

        self._pbar = tk.Canvas(wrap, height=3, bg=C["panel"], highlightthickness=0)
        self._pbar.grid(row=0, column=0, columnspan=2, sticky="ew", padx=14)

        ib = tk.Frame(wrap, bg=C["border2"], padx=1, pady=1)
        ib.grid(row=1, column=0, sticky="ew", padx=10, pady=6); ib.grid_columnconfigure(1, weight=1)
        tk.Label(ib, text=" ❯ ", font=("Consolas",13,"bold"), fg=C["cyan"], bg=C["bg"]).grid(row=0, column=0)
        self._input = tk.Entry(ib, font=("Consolas",12), bg=C["bg"], fg=C["white"],
                               insertbackground=C["cyan"], relief="flat", bd=0)
        self._input.grid(row=0, column=1, sticky="ew", pady=10, padx=(0,10))
        self._input.bind("<Return>", self._send)
        self._input.focus_set()

        self._send_btn = tk.Button(wrap, text="EXECUTE ▶", command=self._send,
                                   bg=C["blue"], fg="white", font=("Consolas",10,"bold"),
                                   relief="flat", bd=0, padx=16, pady=10,
                                   activebackground="#0033aa", cursor="hand2")
        self._send_btn.grid(row=1, column=1, padx=(0,10), pady=6, sticky="ns")

        qf = tk.Frame(wrap, bg=C["panel"])
        qf.grid(row=2, column=0, columnspan=2, sticky="ew", padx=12, pady=(0,4))
        for lbl, cmd in [("📁 파일",self._q_file),("👁 화면",self._q_screen),
                          ("🖼 이미지",self._q_img),("🎬 비디오",self._q_vid),
                          ("🌐 URL",self._q_url),("📧 이메일",self._q_email),
                          ("📸 스크린샷",self._q_shot),("🔄 초기화",self._clear)]:
            tk.Button(qf, text=lbl, command=cmd, bg=C["bg3"], fg=C["mid"],
                      font=("Consolas",8), relief="flat", bd=0, padx=6, pady=3,
                      activebackground=C["border2"], activeforeground=C["cyan"], cursor="hand2").pack(side="left", padx=2)

        self._status_var = tk.StringVar(value="◈  JihunAI v1.1 ONLINE")
        tk.Label(wrap, textvariable=self._status_var, font=("Consolas",8),
                 fg=C["dim"], bg=C["panel"], anchor="e").grid(row=3, column=0, columnspan=2, sticky="ew", padx=16, pady=(0,4))

    def _style_combo(self):
        s = ttk.Style()
        try: s.theme_use("clam")
        except: pass
        s.configure("TCombobox", fieldbackground=C["bg3"], background=C["bg3"],
                     foreground=C["cyan"], arrowcolor=C["cyan"], bordercolor=C["border"])

    # ─────────────── 애니메이션 ───────────────
    def _draw_orb(self):
        c = self._orb; c.delete("all"); cx=cy=20; r=16; t=self._tick
        c.create_oval(cx-r,cy-r,cx+r,cy+r,outline=C["cyan"],width=1)
        for i in range(6):
            a=math.radians(t*3+i*60); px=cx+(r-3)*math.cos(a); py=cy+(r-3)*math.sin(a)
            br=int(60+195*abs(math.sin(a*2)))
            c.create_oval(px-2,py-2,px+2,py+2,fill=f"#{0:02x}{min(255,br):02x}{min(255,br+50):02x}",outline="")
        pr=3+2*abs(math.sin(math.radians(t*5)))
        c.create_oval(cx-pr,cy-pr,cx+pr,cy+pr,fill=C["cyan"],outline="")

    _pbar_x = 0
    _thinking_dots = 0
    _thinking_msgs = ["생각 중","계획 수립 중","도구 실행 중","분석 중","처리 중","검토 중","완성 중","재시도 중"]

    def _animation_loop(self):
        self._tick += 1
        self._draw_orb()
        if self.is_processing:
            self._animate_pbar()
            self._thinking_dots = (self._thinking_dots+1) % 4
            if self._current_task:
                msg = f"⟳ Task {self._current_task['id']}: {self._current_task['title'][:20]}"
            else:
                msg = self._thinking_msgs[(self._tick//20) % len(self._thinking_msgs)]
            self._thinking_lbl.configure(text=f"  {msg}{'.'*self._thinking_dots}")
        self.root.after(40, self._animation_loop)

    def _animate_pbar(self):
        self._pbar.delete("all"); w=self._pbar.winfo_width()
        if w<=1: return
        bw=int(w*0.22); x1=self._pbar_x%(w+bw)-bw
        for i in range(bw):
            alpha=min(1.0,i/20)*min(1.0,(bw-i)/20); br=int(200*alpha)
            self._pbar.create_line(max(0,x1+i),0,max(0,x1+i),3,fill=f"#{0:02x}{br:02x}{min(255,br+55):02x}")
        self._pbar_x += 8

    # ─────────────── 멀티봇 상태 업데이트 ───────────────
    def _update_bot_status(self, active_model=None, success=None, error=None):
        for m_key, dot in self._bot_labels.items():
            stats = self._bot_stats.get(m_key, {})
            if m_key == active_model:
                dot.configure(fg=C["cyan"])
                self._bot_lbl.configure(text=f"[{BOT_ROSTER[m_key][2]}]")
            elif stats.get("errors", 0) > 0 and stats.get("success", 0) == 0:
                dot.configure(fg=C["red"])
            elif stats.get("success", 0) > 0:
                dot.configure(fg=C["green"])
            else:
                dot.configure(fg=C["dim"])

    # ─────────────── 태스크 패널 ───────────────
    def _update_task_panel(self):
        for w in self._task_frame.winfo_children(): w.destroy()
        for task in self._tasks:
            s = task["status"]
            color={"done":C["task_done"],"active":C["task_active"],"pending":C["task_pending"],"fail":C["task_fail"],"retry":C["retry"]}.get(s,C["dim"])
            icon={"done":"✓","active":"⟳","pending":"○","fail":"✗","retry":"↻"}.get(s,"○")
            row=tk.Frame(self._task_frame,bg=C["bg2"]); row.pack(fill="x",pady=1)
            tk.Label(row,text=f"{icon}",fg=color,bg=C["bg2"],font=("Consolas",8)).pack(side="left")
            title=task["title"][:24]+("…" if len(task["title"])>24 else "")
            retries=self._retry_counts.get(task["id"],0)
            retry_str=f" ×{retries}" if retries>0 else ""
            tk.Label(row,text=f" {title}{retry_str}",fg=color if s!="pending" else C["dim"],bg=C["bg2"],font=("Consolas",7)).pack(side="left")

    # ─────────────── 채팅 출력 ───────────────
    def _append(self, role, content, extra=None):
        self._chat.configure(state="normal")
        ts = datetime.now().strftime("%H:%M:%S")
        if role == "user":
            self._chat.insert(tk.END,f"\n  ◈ USER  ","u_lbl"); self._chat.insert(tk.END,f"[{ts}]\n","ts"); self._chat.insert(tk.END,f"  {content}\n","u_txt")
        elif role == "assistant":
            self._chat.insert(tk.END,f"\n  ◈ JihunAI  ","a_lbl"); self._chat.insert(tk.END,f"[{ts}]\n","ts"); self._chat.insert(tk.END,f"  {content}\n","a_txt")
        elif role == "plan":
            self._chat.insert(tk.END,f"\n  ◈ TASK PLAN  ","plan_lbl"); self._chat.insert(tk.END,f"[{ts}]\n","ts"); self._chat.insert(tk.END,f"  {content}\n","plan_txt")
        elif role == "task_start":
            self._chat.insert(tk.END,f"  ⟳ [{extra}]  {content}\n","task_hdr")
        elif role == "task_done":
            self._chat.insert(tk.END,f"  ✓  {content}\n","task_ok")
        elif role == "task_fail":
            self._chat.insert(tk.END,f"  ✗  {content}\n","task_fail")
        elif role == "retry":
            self._chat.insert(tk.END,f"\n  ↻ RETRY  ","retry_lbl"); self._chat.insert(tk.END,f"[{ts}]  {content}\n","ts")
        elif role == "fallback":
            self._chat.insert(tk.END,f"\n  ⟳ FALLBACK  ","fb_lbl"); self._chat.insert(tk.END,f"[{ts}]  {content}\n","ts")
        elif role == "tool_call":
            self._chat.insert(tk.END,f"\n  ⚙ {extra}  ","tc_lbl"); self._chat.insert(tk.END,f"[{ts}]\n","ts"); self._chat.insert(tk.END,f"  {str(content)[:300]}{'...' if len(str(content))>300 else ''}\n","tc_txt")
        elif role == "tool_result":
            self._chat.insert(tk.END,f"  ┗━ {str(content)[:500]}{'...' if len(str(content))>500 else ''}\n","tr_txt")
        elif role == "error":
            self._chat.insert(tk.END,f"\n  ✕ ERROR [{ts}]\n","tc_lbl"); self._chat.insert(tk.END,f"  {content}\n","err_txt")
        elif role == "system":
            self._chat.insert(tk.END,f"\n  ─── {content} ───\n","sys_txt")
        self._chat.configure(state="disabled"); self._chat.see(tk.END)

    # ─────────────── API 호출 (재시도 + fallback) ───────────────
    def _sanitize_messages(self, messages: list) -> list:
        """conversation을 API 전송 전에 완전 정리 — 400 오류 완전 방지

        규칙:
        1. assistant+tool_calls 뒤에는 반드시 그 tool_call_id에 대응하는 tool 메시지들이 와야 함
        2. 대응 tool 메시지가 없는 assistant+tool_calls → tool_calls 제거하고 content만 유지
        3. 대응하는 assistant+tool_calls 없는 tool 메시지 → 제거
        4. 연속된 같은 role → user/user 연속은 합치기
        """
        # 1단계: 기본 정규화
        normalized = []
        for msg in messages:
            role = msg.get("role", "")
            if role == "assistant":
                m = {"role": "assistant", "content": msg.get("content") or ""}
                if msg.get("tool_calls"):
                    m["tool_calls"] = msg["tool_calls"]
                normalized.append(m)
            elif role == "tool":
                normalized.append({
                    "role": "tool",
                    "tool_call_id": msg.get("tool_call_id", ""),
                    "content": str(msg.get("content", ""))[:6000]
                })
            elif role in ("user", "system"):
                normalized.append({"role": role, "content": str(msg.get("content", ""))})

        # 2단계: tool_calls ↔ tool 응답 정합성 검증
        result = []
        i = 0
        while i < len(normalized):
            msg = normalized[i]
            if msg.get("role") == "assistant" and msg.get("tool_calls"):
                needed_ids = [tc["id"] for tc in msg["tool_calls"]]
                # 뒤따라오는 tool 메시지 수집
                j = i + 1
                found = {}
                while j < len(normalized) and normalized[j].get("role") == "tool":
                    tc_id = normalized[j].get("tool_call_id", "")
                    found[tc_id] = normalized[j]
                    j += 1

                all_found = all(nid in found for nid in needed_ids)
                if all_found:
                    # 완전한 쌍 → 그대로 추가
                    result.append(msg)
                    for nid in needed_ids:
                        result.append(found[nid])
                    i = j
                else:
                    # 불완전 → tool_calls 제거하고 content만 보존
                    if msg.get("content"):
                        result.append({"role": "assistant", "content": msg["content"]})
                    # 대응 없는 tool 메시지들도 건너뜀
                    i = j
            elif msg.get("role") == "tool":
                # 앞에 대응하는 assistant+tool_calls 없는 고아 tool → 제거
                i += 1
            else:
                result.append(msg)
                i += 1

        # 3단계: 연속 user 메시지 합치기 (일부 API가 거부함)
        final = []
        for msg in result:
            if final and final[-1].get("role") == "user" and msg.get("role") == "user":
                final[-1]["content"] += "\n" + msg["content"]
            else:
                final.append(msg)

        return final

    def _call_api(self, model_key: str, messages: list, tools=None,
                  tool_choice="auto", max_tokens=4096, temp=0.7,
                  retries=2) -> dict | None:
        """API 호출. 실패 시 retries번 재시도. 반환: choices[0] dict 또는 None"""
        model_id = BOT_ROSTER.get(model_key, (model_key,))[0]
        headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}

        # 메시지 포맷 정리 (400 오류 방지)
        clean_messages = self._sanitize_messages(messages)

        payload = {"model": model_id, "messages": clean_messages,
                   "max_tokens": max_tokens, "temperature": temp}
        if tools:
            payload["tools"] = tools
            payload["tool_choice"] = tool_choice

        for attempt in range(retries + 1):
            try:
                self._bot_stats[model_key]["calls"] = self._bot_stats.get(model_key, {}).get("calls", 0) + 1
                resp = requests.post(API_URL, headers=headers, json=payload, timeout=120)
                if resp.status_code == 200:
                    self._bot_stats[model_key]["success"] = self._bot_stats.get(model_key, {}).get("success", 0) + 1
                    return resp.json()["choices"][0]
                else:
                    err_text = resp.text[:300]
                    self._bot_stats[model_key]["errors"] = self._bot_stats.get(model_key, {}).get("errors", 0) + 1
                    if attempt < retries:
                        self.root.after(0, lambda a=attempt, e=err_text: self._append("retry", f"API {resp.status_code} — 재시도 {a+1}회 ({e[:100]})"))
                        time.sleep(1.5 ** attempt)
                    else:
                        return None
            except requests.exceptions.Timeout:
                self._bot_stats[model_key]["errors"] = self._bot_stats.get(model_key, {}).get("errors", 0) + 1
                if attempt < retries:
                    self.root.after(0, lambda a=attempt: self._append("retry", f"타임아웃 — 재시도 {a+1}회"))
                    time.sleep(2)
                else:
                    return None
            except Exception as e:
                return None
        return None

    # ─────────────── Task Planner ───────────────
    def _call_planner(self, user_msg: str) -> list:
        simple = len(user_msg) < 25 and not any(k in user_msg for k in ["파일","분석","만들","수정","개선","실행","검색","ppt","pptx","PPT","PPTX"])
        if simple: return []
        try:
            choice = self._call_api(
                self._model_var.get(),
                [{"role":"system","content":PLANNER_SYSTEM}, {"role":"user","content":user_msg}],
                tools=None, max_tokens=600, temp=0.3, retries=1
            )
            if not choice: return []
            raw = choice["message"]["content"].strip()
            raw = re.sub(r'^```[a-z]*\n?','',raw); raw=re.sub(r'\n?```$','',raw)
            data = json.loads(raw)
            return [{"id":t["id"],"title":t["title"],"desc":t.get("desc",""),
                     "priority":t.get("priority","medium"),
                     "max_retries":t.get("max_retries",2),
                     "status":"pending"} for t in data.get("tasks",[])]
        except: return []

    # ─────────────── 메인 전송 ───────────────
    def _send(self, event=None):
        if self.is_processing: return
        text = self._input.get().strip()
        if not text: return
        self._input.delete(0, tk.END)
        self.conversation.append({"role":"user","content":text})
        self._append("user", text)
        self._start_proc()
        threading.Thread(target=self._run_agent, args=(text,), daemon=True).start()

    def _send_text(self, text):
        if self.is_processing: return
        self.conversation.append({"role":"user","content":text})
        self._append("user", text)
        self._start_proc()
        threading.Thread(target=self._run_agent, args=(text,), daemon=True).start()

    def _start_proc(self):
        self.is_processing = True; self._pbar_x = 0
        self._tasks = []; self._current_task = None; self._retry_counts = {}
        self._fallback_idx = 0; self._active_model = self._model_var.get()
        self._send_btn.configure(text="⟳ 처리 중...", state="disabled", bg=C["dim"])
        self._status_var.set("◈  AGENT PROCESSING...")
        self.root.after(0, lambda: self._update_bot_status(active_model=self._active_model))

    # ─────────────── ✅ 에이전트 루프 (v1.1) ───────────────
    def _run_agent(self, user_msg: str):
        # 1. Task Planner
        self.root.after(0, lambda: self._status_var.set("◈  PLANNING..."))
        tasks = self._call_planner(user_msg)
        if tasks:
            self._tasks = tasks
            plan_lines = "\n".join(f"  {t['id']}. [{t['priority'].upper()}] {t['title']} — {t['desc']}" for t in tasks)
            self.root.after(0, lambda: self._append("plan", f"총 {len(tasks)}개 태스크\n{plan_lines}"))
            self.root.after(0, self._update_task_panel)

        # 2. 에이전트 루프
        current_model = self._active_model
        max_iter = BOT_ROSTER.get(current_model, (None, 16))[1]
        consecutive_empty = 0  # 도구 호출도 없고 내용도 없는 응답 연속 횟수
        iteration = 0

        def _try_next_fallback(reason: str) -> bool:
            """다음 fallback 모델로 전환. 성공하면 True"""
            nonlocal current_model, max_iter
            self._fallback_idx += 1
            if self._fallback_idx >= len(FALLBACK_CHAIN):
                return False
            next_model = FALLBACK_CHAIN[self._fallback_idx]
            if next_model == current_model: self._fallback_idx += 1
            if self._fallback_idx >= len(FALLBACK_CHAIN): return False
            next_model = FALLBACK_CHAIN[self._fallback_idx]
            current_model = next_model
            max_iter = BOT_ROSTER.get(current_model, (None, 14))[1]
            # fallback 시 conversation 정리: 미완성 tool_calls 메시지 제거
            # (assistant tool_calls 이후 tool 응답이 없는 경우 제거)
            cleaned = []
            i = 0
            msgs = self.conversation
            while i < len(msgs):
                m = msgs[i]
                if m.get("role") == "assistant" and m.get("tool_calls"):
                    # 뒤에 tool 응답이 모두 있는지 확인
                    needed_ids = {tc["id"] for tc in m["tool_calls"]}
                    found_ids = set()
                    j = i + 1
                    while j < len(msgs) and msgs[j].get("role") == "tool":
                        found_ids.add(msgs[j].get("tool_call_id",""))
                        j += 1
                    if needed_ids == found_ids:
                        cleaned.extend(msgs[i:j])
                        i = j
                    else:
                        # 미완성 tool 교환 → 전부 제거 (400 오류 원인)
                        i = j
                else:
                    cleaned.append(m)
                    i += 1
            self.conversation = cleaned
            self.root.after(0, lambda m=next_model, r=reason: (
                self._append("fallback", f"{r} → fallback: {BOT_ROSTER.get(m,(m,))[2]}"),
                self._update_bot_status(active_model=m)
            ))
            return True

        try:
            while iteration < max_iter:
                iteration += 1
                self.root.after(0, lambda i=iteration, m=current_model: self._status_var.set(f"◈  [{m.split('-')[0].upper()}] iter {i}"))

                choice = self._call_api(current_model, self.conversation, tools=TOOLS,
                                        tool_choice="auto", max_tokens=4096, retries=2)

                # API 완전 실패 → fallback
                if choice is None:
                    if not _try_next_fallback("API 호출 실패"):
                        self.root.after(0, lambda: self._append("error", "모든 봇 호출 실패. 나중에 다시 시도해주세요."))
                        break
                    iteration = 0  # fallback 후 이터레이션 리셋
                    continue

                msg    = choice["message"]
                finish = choice.get("finish_reason", "stop")
                content = msg.get("content") or ""

                # ── 태그 파싱 ──
                def parse_tags(text):
                    # TASK_PLAN
                    for m in re.finditer(r'<TASK_PLAN>(.*?)</TASK_PLAN>', text, re.DOTALL):
                        try:
                            d = json.loads(m.group(1))
                            if not self._tasks:
                                self._tasks = [{"id":t["id"],"title":t["title"],"desc":t.get("desc",""),"priority":t.get("priority","medium"),"max_retries":t.get("max_retries",2),"status":"pending"} for t in d.get("tasks",[])]
                                plan_lines = "\n".join(f"  {t['id']}. {t['title']}" for t in self._tasks)
                                self.root.after(0, lambda p=plan_lines: self._append("plan", f"총 {len(self._tasks)}개 태스크\n{p}"))
                                self.root.after(0, self._update_task_panel)
                        except: pass
                    # TASK_START
                    for m in re.finditer(r'<TASK_START>(.*?)</TASK_START>', text, re.DOTALL):
                        try:
                            d=json.loads(m.group(1)); tid=d.get("id")
                            for t in self._tasks:
                                if t["id"]==tid:
                                    t["status"]="active"; self._current_task=t
                                    self.root.after(0, lambda tt=t: (self._append("task_start",tt["title"],extra=f"Task {tt['id']}"), self._update_task_panel()))
                                    break
                        except: pass
                    # TASK_DONE
                    for m in re.finditer(r'<TASK_DONE>(.*?)</TASK_DONE>', text, re.DOTALL):
                        try:
                            d=json.loads(m.group(1)); tid=d.get("id"); res=d.get("result","완료")
                            for t in self._tasks:
                                if t["id"]==tid:
                                    t["status"]="done"; self._current_task=None
                                    self.root.after(0, lambda tt=t, r=res: (self._append("task_done",f"{tt['title']}: {r}"), self._update_task_panel()))
                                    break
                        except: pass
                    # TASK_FAIL
                    for m in re.finditer(r'<TASK_FAIL>(.*?)</TASK_FAIL>', text, re.DOTALL):
                        try:
                            d=json.loads(m.group(1)); tid=d.get("id"); reason=d.get("reason","실패")
                            for t in self._tasks:
                                if t["id"]==tid:
                                    self._retry_counts[tid] = self._retry_counts.get(tid,0)+1
                                    max_r = t.get("max_retries",2)
                                    if self._retry_counts[tid] <= max_r:
                                        t["status"]="retry"
                                        self.root.after(0, lambda tt=t, rc=self._retry_counts[tid], r=reason: (self._append("retry",f"Task {tt['id']} 재시도 {rc}회: {r}"), self._update_task_panel()))
                                    else:
                                        t["status"]="fail"; self._current_task=None
                                        self.root.after(0, lambda tt=t, r=reason: (self._append("task_fail",f"{tt['title']}: {r}"), self._update_task_panel()))
                                    break
                        except: pass
                    # 태그 제거
                    clean = re.sub(r'<TASK_(?:PLAN|START|DONE|FAIL)>.*?</TASK_(?:PLAN|START|DONE|FAIL)>','',text,flags=re.DOTALL)
                    return clean.strip()

                # ── tool_calls ──
                if finish == "tool_calls" and msg.get("tool_calls"):
                    consecutive_empty = 0

                    # assistant 메시지를 tool_calls 포함해서 conversation에 추가
                    # API 호환성: content가 None이면 빈 문자열로
                    assistant_msg = {
                        "role": "assistant",
                        "content": content or "",
                        "tool_calls": msg["tool_calls"]
                    }
                    self.conversation.append(assistant_msg)

                    if content:
                        clean = parse_tags(content)
                        if clean: self.root.after(0, lambda c=clean: self._append("assistant", c))

                    # 각 tool 실행 후 즉시 tool 메시지 추가 (순서 보장)
                    for tc in msg["tool_calls"]:
                        fn = tc["function"]["name"]
                        try: fn_args = json.loads(tc["function"]["arguments"])
                        except: fn_args = {}
                        args_str = json.dumps(fn_args, ensure_ascii=False)[:250]
                        self.root.after(0, lambda n=fn, a=args_str: self._append("tool_call", a, extra=n))
                        self.root.after(0, lambda n=fn: self._status_var.set(f"◈  EXEC: {n}..."))

                        def _st(m, n=fn): self.root.after(0, lambda: self._status_var.set(f"◈  {n}: {m}"))
                        result = execute_tool(fn, fn_args, st=_st)
                        result_str = str(result)[:6000]
                        self.root.after(0, lambda r=result_str: self._append("tool_result", r))

                        # 각 tool 결과를 바로 conversation에 추가 (assistant 메시지 바로 뒤에 위치)
                        self.conversation.append({
                            "role": "tool",
                            "tool_call_id": tc["id"],
                            "content": result_str
                        })

                    continue

                else:
                    # 일반 응답
                    if content:
                        consecutive_empty = 0
                        clean = parse_tags(content)
                        if clean:
                            self.conversation.append({"role":"assistant","content":content})
                            self.root.after(0, lambda c=clean: self._append("assistant", c))
                    else:
                        # 내용 없는 응답 연속 시 fallback
                        consecutive_empty += 1
                        if consecutive_empty >= 2:
                            if not _try_next_fallback(f"빈 응답 {consecutive_empty}회 연속"):
                                break
                            consecutive_empty = 0; iteration = 0
                            continue

                    # ── 루프 종료 판단 ──
                    if self._tasks:
                        pending = [t for t in self._tasks if t["status"] in ("pending","active","retry")]
                        all_done = len(pending) == 0
                        if all_done:
                            # 모든 태스크 완료 → 즉시 종료
                            self.root.after(0, self._update_task_panel)
                            break
                        elif content and pending:
                            # 아직 남은 태스크 있고 AI가 내용 생성 중 → 계속
                            continue
                        else:
                            # 태스크 있는데 AI가 아무것도 안 함 → 종료
                            for t in self._tasks:
                                if t["status"] in ("active","pending"):
                                    t["status"] = "done"
                            self.root.after(0, self._update_task_panel)
                            break
                    else:
                        # 태스크 없는 단순 대화 → finish_reason == stop 이면 종료
                        if finish == "stop" or not content:
                            break
                        # content 있고 finish != stop 이면 계속 (스트리밍 미완성 대비)
                        break  # 단순 대화는 1회 응답으로 종료

        except Exception as e:
            import traceback
            self.root.after(0, lambda e=traceback.format_exc()[-700:]: self._append("error", e))
        finally:
            self.root.after(0, self._done_proc)
            self.root.after(0, self._save_chat)

    def _done_proc(self):
        self.is_processing = False
        self._thinking_lbl.configure(text="")
        self._send_btn.configure(text="EXECUTE ▶", state="normal", bg=C["blue"])
        self._status_var.set("◈  JihunAI v1.4 ONLINE  —  READY")
        self._pbar.delete("all"); self._input.focus_set()
        self._update_bot_status()

    # ─────────────── 빠른 액션 ───────────────
    def _q_file(self):
        p=filedialog.askopenfilename(filetypes=[("모든 파일","*.*")])
        if p:
            ext=Path(p).suffix.lower()
            if ext in [".png",".jpg",".jpeg",".bmp",".gif",".webp"]: self._send_text(f"이 이미지를 Vision AI로 분석해줘: {p}")
            elif ext in [".mp4",".avi",".mov",".mkv"]: self._send_text(f"이 비디오를 분석해줘: {p}")
            elif ext in [".pptx",".ppt"]: self._send_text(f"이 PPT 파일을 훨씬 예쁘게 개선해줘: {p}")
            else: self._input.insert(tk.END, p)
    def _q_screen(self): self._send_text("내 화면을 캡처해서 Vision AI로 분석해줘.")
    def _q_img(self):
        p=filedialog.askopenfilename(filetypes=[("이미지","*.png *.jpg *.jpeg *.bmp *.gif *.webp")])
        if p: self._send_text(f"이 이미지를 Vision AI로 분석해줘: {p}")
    def _q_vid(self):
        p=filedialog.askopenfilename(filetypes=[("비디오","*.mp4 *.avi *.mov *.mkv")])
        if p: self._send_text(f"이 비디오를 5프레임으로 분석해줘: {p}")
    def _q_url(self):
        import tkinter.simpledialog
        u=tkinter.simpledialog.askstring("URL","URL:")
        if u:
            if not u.startswith("http"): u="https://"+u
            self._send_text(f"이 URL 열어줘: {u}")
    def _q_email(self):
        self._input.delete(0,tk.END); self._input.insert(0,"Gmail 작성 창 열어줘. 수신: , 제목: , 내용: ")
    def _q_shot(self): self._send_text("화면 캡처해서 screenshot.png로 저장해줘.")
    def _clear(self):
        if messagebox.askyesno("초기화","대화를 초기화할까요?"):
            self.conversation=[{"role":"system","content":SYSTEM_PROMPT}]
            self._chat.configure(state="normal"); self._chat.delete("1.0",tk.END); self._chat.configure(state="disabled")
            self._tasks=[]; self._update_task_panel(); self._append("system","CLEARED — JihunAI v1.1 READY")

    def _config(self):
        cfg=tk.Toplevel(self.root); cfg.title("JihunAI 설정"); cfg.geometry("460x320"); cfg.configure(bg=C["panel"])
        tk.Label(cfg,text="◈  EMAIL CONFIG",font=("Consolas",11,"bold"),fg=C["cyan"],bg=C["panel"]).pack(pady=14)
        entries={}
        for lbl,key in [("이메일","email"),("앱 비밀번호","password"),("SMTP 호스트","smtp_host"),("SMTP 포트","smtp_port")]:
            r=tk.Frame(cfg,bg=C["panel"]); r.pack(fill="x",padx=30,pady=3)
            tk.Label(r,text=f"{lbl}:",font=("Consolas",9),fg=C["mid"],bg=C["panel"],width=14,anchor="e").pack(side="left")
            e=tk.Entry(r,font=("Consolas",10),bg=C["bg3"],fg=C["white"],insertbackground=C["cyan"],relief="flat",show="*" if key=="password" else "")
            e.insert(0,str(EMAIL_CONFIG.get(key,""))); e.pack(side="left",fill="x",expand=True,padx=(8,0))
            entries[key]=e
        def save():
            for k,e in entries.items(): EMAIL_CONFIG[k]=int(e.get()) if k=="smtp_port" else e.get()
            messagebox.showinfo("저장","저장 완료"); cfg.destroy()
        tk.Button(cfg,text="저장",command=save,bg=C["blue"],fg="white",font=("Consolas",10,"bold"),relief="flat",padx=18,pady=6,cursor="hand2").pack(pady=16)
        tk.Label(cfg,text="* pip install python-pptx 설치 시 PPTX 직접 편집 가능",font=("Consolas",7),fg=C["dim"],bg=C["panel"]).pack()

    # ─────────────── 채팅 저장/로드 ───────────────
    def _save_chat(self):
        msgs=[m for m in self.conversation if m.get("role")!="system"]
        if not msgs: return
        if not self.current_chat_id: self.current_chat_id=datetime.now().strftime("%Y%m%d_%H%M%S")
        path=os.path.join(CHATS_DIR,f"{self.current_chat_id}.json")
        save=[]
        for m in self.conversation:
            c=m.get("content","")
            if isinstance(c,list): c=" ".join(b.get("text","") for b in c if b.get("type")=="text")
            save.append({**m,"content":c})
        with open(path,"w",encoding="utf-8") as f: json.dump(save,f,ensure_ascii=False,indent=2)
        self.load_chat_list()

    def load_chat_list(self):
        self._chat_lb.delete(0,tk.END)
        files=sorted(glob.glob(os.path.join(CHATS_DIR,"*.json")),reverse=True)
        for file in files:
            cid=os.path.basename(file).replace(".json","")
            try:
                with open(file,encoding="utf-8") as f: data=json.load(f)
                prev=next((str(m.get("content",""))[:20] for m in data if m.get("role")=="user"),cid)
            except: prev=cid
            self._chat_lb.insert(tk.END,f"  {cid[:10]}  {prev}...")
        self._chat_files=files

    def load_latest_chat(self):
        files=sorted(glob.glob(os.path.join(CHATS_DIR,"*.json")),reverse=True)
        if files: self._load_file(files[0])
        else:
            self._append("system","JihunAI v1.1 — 지훈이의 AI 에이전트")
            self._append("system","Multi-Bot Fallback · Task Planner · Vision AI · PPTX 직접 편집")

    def _load_file(self,path):
        try:
            with open(path,encoding="utf-8") as f: self.conversation=json.load(f)
            self.current_chat_id=os.path.basename(path).replace(".json","")
            self._chat.configure(state="normal"); self._chat.delete("1.0",tk.END); self._chat.configure(state="disabled")
            for m in self.conversation:
                role,c=m.get("role",""),m.get("content","")
                if isinstance(c,list): c=" ".join(b.get("text","") for b in c if b.get("type")=="text")
                if role=="user" and c: self._append("user",c)
                elif role=="assistant" and c: self._append("assistant",c)
        except Exception as e: self._append("error",f"로드 실패: {e}")

    def _on_chat_sel(self,e):
        sel=self._chat_lb.curselection()
        if sel and hasattr(self,"_chat_files") and sel[0]<len(self._chat_files):
            self._load_file(self._chat_files[sel[0]])

    def _new_chat(self):
        self._save_chat()
        self.conversation=[{"role":"system","content":SYSTEM_PROMPT}]
        self.current_chat_id=None
        self._chat.configure(state="normal"); self._chat.delete("1.0",tk.END); self._chat.configure(state="disabled")
        self._tasks=[]; self._update_task_panel()
        self._append("system","NEW MISSION — JihunAI v1.1 READY")
        self.load_chat_list()

    def _del_chat(self):
        if not self.current_chat_id: return
        if messagebox.askyesno("삭제","삭제할까요?"):
            p=os.path.join(CHATS_DIR,f"{self.current_chat_id}.json")
            if os.path.exists(p): os.remove(p)
            self._new_chat()

    def _on_close(self): self._save_chat(); self.root.destroy()


# ──────────────────────── 진입점 ────────────────────────
if __name__ == "__main__":
    print("╔══════════════════════════════════════════════════════╗")
    print("║    JihunAI v1.1 — Multi-Bot · Task Planner          ║")
    print("╠══════════════════════════════════════════════════════╣")
    for nm, key in [("PIL/Vision","PIL"),("python-pptx","pptx"),("PyAutoGUI","pyautogui"),("OpenCV","cv2"),("Psutil","psutil")]:
        status = "✓" if LIBS.get(key) else f"✗  →  pip install {key.lower().replace('_','-')}"
        print(f"║  {nm:<16}: {status:<36}║")
    print("╚══════════════════════════════════════════════════════╝")
    if not LIBS.get("pptx"):
        print("\n⚠  PPTX 직접 편집을 위해:  pip install python-pptx")
    root = tk.Tk()
    app = JihunAI(root)
    root.mainloop()
